#!/usr/bin/env python3
"""
PPTX to HTML Converter
Converts PowerPoint presentations to interactive HTML format
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import base64
from io import BytesIO

def extract_image_as_base64(image):
    """Extract image and convert to base64 data URI"""
    try:
        image_stream = BytesIO(image.blob)
        image_bytes = image_stream.getvalue()
        image_base64 = base64.b64encode(image_bytes).decode('utf-8')

        # Detect image format
        ext = image.ext if hasattr(image, 'ext') else 'png'
        mime_type = f'image/{ext}'

        return f'data:{mime_type};base64,{image_base64}'
    except Exception as e:
        print(f"Error extracting image: {e}")
        return None

def extract_text_from_shape(shape):
    """Extract text from a shape"""
    if hasattr(shape, "text"):
        return shape.text
    return ""

def convert_slide_to_html(slide, slide_num, total_slides):
    """Convert a single slide to HTML"""
    html = f'<div class="slide" id="slide-{slide_num}" data-slide="{slide_num}">\n'
    html += f'  <div class="slide-content">\n'
    html += f'    <div class="slide-number">{slide_num}/{total_slides}</div>\n'

    # Process shapes in the slide
    for shape in slide.shapes:
        # Handle text boxes
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()

            # Determine if it's a title or content
            if shape.is_placeholder:
                placeholder = shape.placeholder_format
                if placeholder.type == 1:  # Title
                    html += f'    <h1 class="slide-title">{text}</h1>\n'
                elif placeholder.type == 2:  # Body/Content
                    # Split into lines and create list
                    lines = text.split('\n')
                    if len(lines) > 1:
                        html += '    <ul class="slide-content-list">\n'
                        for line in lines:
                            if line.strip():
                                html += f'      <li>{line.strip()}</li>\n'
                        html += '    </ul>\n'
                    else:
                        html += f'    <p class="slide-text">{text}</p>\n'
            else:
                html += f'    <p class="slide-text">{text}</p>\n'

        # Handle images
        if hasattr(shape, "image"):
            image_data = extract_image_as_base64(shape.image)
            if image_data:
                html += f'    <img src="{image_data}" class="slide-image" alt="Slide {slide_num} Image">\n'

        # Handle tables
        if hasattr(shape, "table"):
            html += '    <table class="slide-table">\n'
            table = shape.table
            for row in table.rows:
                html += '      <tr>\n'
                for cell in row.cells:
                    html += f'        <td>{cell.text}</td>\n'
                html += '      </tr>\n'
            html += '    </table>\n'

    html += '  </div>\n'
    html += '</div>\n'

    return html

def generate_html_template(slides_html, title="PowerPoint Presentation"):
    """Generate complete HTML document"""
    html = f"""<!DOCTYPE html>
<html lang="tr">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}

        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            display: flex;
            justify-content: center;
            align-items: center;
            padding: 20px;
        }}

        .presentation-container {{
            width: 100%;
            max-width: 1200px;
            background: white;
            border-radius: 20px;
            box-shadow: 0 20px 60px rgba(0,0,0,0.3);
            overflow: hidden;
        }}

        .slide {{
            display: none;
            padding: 60px;
            min-height: 600px;
            position: relative;
            animation: slideIn 0.5s ease-in-out;
        }}

        .slide.active {{
            display: block;
        }}

        @keyframes slideIn {{
            from {{
                opacity: 0;
                transform: translateX(50px);
            }}
            to {{
                opacity: 1;
                transform: translateX(0);
            }}
        }}

        .slide-number {{
            position: absolute;
            top: 20px;
            right: 20px;
            background: rgba(0,0,0,0.1);
            padding: 8px 16px;
            border-radius: 20px;
            font-size: 14px;
            color: #666;
            font-weight: 600;
        }}

        .slide-title {{
            font-size: 48px;
            color: #667eea;
            margin-bottom: 30px;
            font-weight: 700;
            line-height: 1.2;
        }}

        .slide-text {{
            font-size: 24px;
            color: #333;
            margin-bottom: 20px;
            line-height: 1.6;
        }}

        .slide-content-list {{
            list-style: none;
            margin-top: 30px;
        }}

        .slide-content-list li {{
            font-size: 28px;
            color: #444;
            margin-bottom: 20px;
            padding-left: 40px;
            position: relative;
            line-height: 1.5;
        }}

        .slide-content-list li:before {{
            content: "→";
            position: absolute;
            left: 0;
            color: #667eea;
            font-weight: bold;
            font-size: 32px;
        }}

        .slide-image {{
            max-width: 100%;
            max-height: 400px;
            margin: 20px auto;
            display: block;
            border-radius: 10px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.2);
        }}

        .slide-table {{
            width: 100%;
            margin: 20px 0;
            border-collapse: collapse;
        }}

        .slide-table td {{
            padding: 15px;
            border: 2px solid #667eea;
            font-size: 18px;
        }}

        .slide-table tr:first-child td {{
            background: #667eea;
            color: white;
            font-weight: bold;
        }}

        .controls {{
            display: flex;
            justify-content: space-between;
            align-items: center;
            padding: 30px 60px;
            background: linear-gradient(to right, #667eea, #764ba2);
        }}

        .nav-button {{
            background: white;
            border: none;
            padding: 15px 30px;
            border-radius: 50px;
            font-size: 18px;
            font-weight: 600;
            color: #667eea;
            cursor: pointer;
            transition: all 0.3s;
            box-shadow: 0 5px 15px rgba(0,0,0,0.2);
        }}

        .nav-button:hover:not(:disabled) {{
            transform: translateY(-2px);
            box-shadow: 0 8px 20px rgba(0,0,0,0.3);
        }}

        .nav-button:disabled {{
            opacity: 0.3;
            cursor: not-allowed;
        }}

        .progress-bar {{
            flex: 1;
            margin: 0 30px;
            height: 8px;
            background: rgba(255,255,255,0.3);
            border-radius: 10px;
            overflow: hidden;
        }}

        .progress-fill {{
            height: 100%;
            background: white;
            border-radius: 10px;
            transition: width 0.3s ease;
        }}

        .fullscreen-button {{
            background: white;
            border: none;
            padding: 12px 20px;
            border-radius: 50px;
            font-size: 16px;
            font-weight: 600;
            color: #667eea;
            cursor: pointer;
            transition: all 0.3s;
            margin-left: 15px;
        }}

        .fullscreen-button:hover {{
            transform: scale(1.05);
        }}

        /* Keyboard shortcuts hint */
        .shortcuts {{
            position: fixed;
            bottom: 20px;
            left: 20px;
            background: rgba(0,0,0,0.7);
            color: white;
            padding: 15px 20px;
            border-radius: 10px;
            font-size: 12px;
            opacity: 0;
            transition: opacity 0.3s;
        }}

        .shortcuts.show {{
            opacity: 1;
        }}

        @media (max-width: 768px) {{
            .slide {{
                padding: 30px;
                min-height: 400px;
            }}

            .slide-title {{
                font-size: 32px;
            }}

            .slide-text {{
                font-size: 18px;
            }}

            .slide-content-list li {{
                font-size: 20px;
            }}

            .controls {{
                padding: 20px 30px;
            }}

            .nav-button {{
                padding: 10px 20px;
                font-size: 14px;
            }}
        }}
    </style>
</head>
<body>
    <div class="presentation-container">
        {slides_html}

        <div class="controls">
            <button class="nav-button" id="prevBtn" onclick="changeSlide(-1)">← Önceki</button>
            <div class="progress-bar">
                <div class="progress-fill" id="progressBar"></div>
            </div>
            <button class="nav-button" id="nextBtn" onclick="changeSlide(1)">Sonraki →</button>
            <button class="fullscreen-button" onclick="toggleFullscreen()">⛶ Tam Ekran</button>
        </div>
    </div>

    <div class="shortcuts" id="shortcuts">
        <strong>Klavye Kısayolları:</strong><br>
        ← → : Slayt değiştir | F: Tam ekran | Esc: Çıkış
    </div>

    <script>
        let currentSlide = 1;
        const totalSlides = document.querySelectorAll('.slide').length;

        // Show first slide
        showSlide(currentSlide);

        // Show shortcuts hint for 3 seconds
        setTimeout(() => {{
            document.getElementById('shortcuts').classList.add('show');
            setTimeout(() => {{
                document.getElementById('shortcuts').classList.remove('show');
            }}, 3000);
        }}, 500);

        function showSlide(n) {{
            const slides = document.querySelectorAll('.slide');

            if (n > totalSlides) currentSlide = totalSlides;
            if (n < 1) currentSlide = 1;

            slides.forEach(slide => slide.classList.remove('active'));
            slides[currentSlide - 1].classList.add('active');

            // Update progress bar
            const progress = (currentSlide / totalSlides) * 100;
            document.getElementById('progressBar').style.width = progress + '%';

            // Update button states
            document.getElementById('prevBtn').disabled = currentSlide === 1;
            document.getElementById('nextBtn').disabled = currentSlide === totalSlides;
        }}

        function changeSlide(direction) {{
            currentSlide += direction;
            showSlide(currentSlide);
        }}

        function toggleFullscreen() {{
            if (!document.fullscreenElement) {{
                document.documentElement.requestFullscreen();
            }} else {{
                document.exitFullscreen();
            }}
        }}

        // Keyboard navigation
        document.addEventListener('keydown', (e) => {{
            if (e.key === 'ArrowLeft') {{
                changeSlide(-1);
            }} else if (e.key === 'ArrowRight' || e.key === ' ') {{
                e.preventDefault();
                changeSlide(1);
            }} else if (e.key === 'f' || e.key === 'F') {{
                toggleFullscreen();
            }}
        }});

        // Touch/swipe support
        let touchStartX = 0;
        let touchEndX = 0;

        document.addEventListener('touchstart', (e) => {{
            touchStartX = e.changedTouches[0].screenX;
        }});

        document.addEventListener('touchend', (e) => {{
            touchEndX = e.changedTouches[0].screenX;
            handleSwipe();
        }});

        function handleSwipe() {{
            if (touchEndX < touchStartX - 50) {{
                changeSlide(1);
            }}
            if (touchEndX > touchStartX + 50) {{
                changeSlide(-1);
            }}
        }}
    </script>
</body>
</html>"""

    return html

def convert_pptx_to_html(pptx_path, output_path=None):
    """Main conversion function"""
    print(f"📂 Reading PowerPoint: {pptx_path}")

    try:
        # Load presentation
        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)

        print(f"📊 Found {total_slides} slides")

        # Convert each slide
        slides_html = ""
        for idx, slide in enumerate(prs.slides, start=1):
            print(f"  Converting slide {idx}/{total_slides}...")
            slide_html = convert_slide_to_html(slide, idx, total_slides)
            slides_html += slide_html

        # Generate complete HTML
        title = Path(pptx_path).stem
        html_content = generate_html_template(slides_html, title)

        # Determine output path
        if output_path is None:
            output_path = Path(pptx_path).with_suffix('.html')

        # Write HTML file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)

        print(f"✅ HTML generated: {output_path}")
        print(f"📏 File size: {os.path.getsize(output_path) / 1024:.2f} KB")

        return output_path

    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return None

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python pptx_to_html_converter.py <pptx_file> [output_html]")
        print("\nExample:")
        print("  python pptx_to_html_converter.py presentation.pptx")
        print("  python pptx_to_html_converter.py presentation.pptx output.html")
        sys.exit(1)

    pptx_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None

    if not os.path.exists(pptx_file):
        print(f"❌ File not found: {pptx_file}")
        sys.exit(1)

    convert_pptx_to_html(pptx_file, output_file)
