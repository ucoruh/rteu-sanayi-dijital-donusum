"""
SANAYİDE DİJİTAL DÖNÜŞÜM - ULTİMATE 2 SAATLİK SUNUM
Dr. Öğr. Üyesi Uğur CORUH
RTEÜ Erasmus+ Programı
19 Ekim 2025 | 10:00-12:00

✨ ÖZELLİKLER:
- İki sunumun harmanlanması (final + full_workshop)
- DETAYLI proje yazım formatı ve yönergeleri
- Her slayt tam ekran Unsplash görselli
- Adım adım öğrenci rehberi
- TÜBİTAK/KOSGEB format şablonları
- 8 modern araç detaylı kullanım

ZAMANLAMA:
- 10:00-10:10 (10dk): Açılış
- 10:10-10:30 (20dk): Sanayi 4.0
- 10:30-10:55 (25dk): Modern Araçlar + DEMO
- 10:55-11:20 (25dk): Fonlama
- 11:20-11:30 (10dk): Mola
- 11:30-12:00 (30dk): Atölye (detaylı formatlarla)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
from PIL import Image
import numpy as np
from PIL import Image, ImageEnhance, ImageFilter

class UltimatePresentation:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        self.stats = {'total': 0, 'success': 0, 'failed': 0}

        # Parlak, canlı renkler
        self.colors = {
            'primary': RGBColor(0, 123, 255),
            'secondary': RGBColor(255, 107, 0),
            'success': RGBColor(40, 199, 111),
            'warning': RGBColor(255, 193, 7),
            'danger': RGBColor(220, 53, 69),
            'purple': RGBColor(130, 88, 255),
            'teal': RGBColor(32, 201, 151),
            'pink': RGBColor(255, 71, 133),
            'dark': RGBColor(33, 37, 41),
            'white': RGBColor(255, 255, 255),
            'light': RGBColor(248, 249, 250),
        }

        # Zengin Unsplash görsel kütüphanesi (57 görsel)
        self.images = {
            # Mevcut görseller
            'hero': 'https://images.unsplash.com/photo-1581091226825-a6a2a5aee158?w=1920&h=1080&fit=crop',
            'factory': 'https://images.unsplash.com/photo-1565043666747-69f6646db940?w=1920&h=1080&fit=crop',
            'technology': 'https://images.unsplash.com/photo-1451187580459-43490279c0fa?w=1920&h=1080&fit=crop',
            'ai_brain': 'https://images.unsplash.com/photo-1677442136019-21780ecad995?w=1920&h=1080&fit=crop',
            'code': 'https://images.unsplash.com/photo-1461749280684-dccba630e2f6?w=1920&h=1080&fit=crop',
            'team': 'https://images.unsplash.com/photo-1522071820081-009f0129c71c?w=1920&h=1080&fit=crop',
            'innovation': 'https://images.unsplash.com/photo-1485827404703-89b55fcc595e?w=1920&h=1080&fit=crop',
            'data': 'https://images.unsplash.com/photo-1551288049-bebda4e38f71?w=1920&h=1080&fit=crop',
            'laptop': 'https://images.unsplash.com/photo-1484788984921-03950022c9ef?w=1920&h=1080&fit=crop',
            'workspace': 'https://images.unsplash.com/photo-1497366216548-37526070297c?w=1920&h=1080&fit=crop',
            'brainstorm': 'https://images.unsplash.com/photo-1559136555-9303baea8ebd?w=1920&h=1080&fit=crop',
            'meeting': 'https://images.unsplash.com/photo-1517245386807-bb43f82c33c4?w=1920&h=1080&fit=crop',
            'success': 'https://images.unsplash.com/photo-1552581234-26160f608093?w=1920&h=1080&fit=crop',
            'growth': 'https://images.unsplash.com/photo-1460925895917-afdab827c52f?w=1920&h=1080&fit=crop',
            'money': 'https://images.unsplash.com/photo-1579621970563-ebec7560ff3e?w=1920&h=1080&fit=crop',
            'presentation': 'https://images.unsplash.com/photo-1558618666-fcd25c85cd64?w=1920&h=1080&fit=crop',
            'network': 'https://images.unsplash.com/photo-1558494949-ef010cbdcc31?w=1920&h=1080&fit=crop',
            'rocket': 'https://images.unsplash.com/photo-1516849841032-87cbac4d88f7?w=1920&h=1080&fit=crop',
            'future': 'https://images.unsplash.com/photo-1451187580459-43490279c0fa?w=1920&h=1080&fit=crop',
            'typing': 'https://images.unsplash.com/photo-1484480974693-6ca0a78fb36b?w=1920&h=1080&fit=crop',
            'notebook': 'https://images.unsplash.com/photo-1455390582262-044cdead277a?w=1920&h=1080&fit=crop',
            'document': 'https://images.unsplash.com/photo-1507842217343-583bb7270b66?w=1920&h=1080&fit=crop',
            'checklist': 'https://images.unsplash.com/photo-1484480974693-6ca0a78fb36b?w=1920&h=1080&fit=crop',
            'planning': 'https://images.unsplash.com/photo-1531498860502-7c67cf02f657?w=1920&h=1080&fit=crop',
            'collaboration': 'https://images.unsplash.com/photo-1553877522-43269d4ea984?w=1920&h=1080&fit=crop',
            'startup': 'https://images.unsplash.com/photo-1519834785169-98be25ec3f84?w=1920&h=1080&fit=crop',
            'writing': 'https://images.unsplash.com/photo-1455390582262-044cdead277a?w=1920&h=1080&fit=crop',
            'tools': 'https://images.unsplash.com/photo-1581091226825-a6a2a5aee158?w=1920&h=1080&fit=crop',
            'hologram': 'https://images.unsplash.com/photo-1550751827-4bd374c3f58b?w=1920&h=1080&fit=crop',

            # Yeni görseller - Açılış bölümü
            'learning': 'https://images.unsplash.com/photo-1523580494863-6f3031224c94?w=1920&h=1080&fit=crop',
            'target_goals': 'https://images.unsplash.com/photo-1519389950473-47ba0277781c?w=1920&h=1080&fit=crop',
            'expectations': 'https://images.unsplash.com/photo-1557804506-669a67965ba0?w=1920&h=1080&fit=crop',

            # Design Thinking - Lean - Agile
            'design_thinking': 'https://raw.githubusercontent.com/ucoruh/rteu-sanayi-dijital-donusum/main/design-thinking-lean-agile-slide1.png',

            # Yeni görseller - Sanayi 4.0
            'industry_intro': 'https://images.unsplash.com/photo-1581091226825-a6a2a5aee158?w=1920&h=1080&fit=crop',
            'revolution': 'https://images.unsplash.com/photo-1518770660439-4636190af475?w=1920&h=1080&fit=crop',
            'iot_sensors': 'https://images.unsplash.com/photo-1558346490-a72e53ae2d4f?w=1920&h=1080&fit=crop',
            'artificial_intelligence': 'https://images.unsplash.com/photo-1555949963-aa79dcee981c?w=1920&h=1080&fit=crop',
            'big_data': 'https://images.unsplash.com/photo-1527474305487-b87b222841cc?w=1920&h=1080&fit=crop',
            'cloud_computing': 'https://images.unsplash.com/photo-1544197150-b99a580bb7a8?w=1920&h=1080&fit=crop',
            'digital_twin': 'https://images.unsplash.com/photo-1620712943543-bcc4688e7485?w=1920&h=1080&fit=crop',
            'robotics': 'https://images.unsplash.com/photo-1485827404703-89b55fcc595e?w=1920&h=1080&fit=crop',
            'turkey_industry': 'https://images.unsplash.com/photo-1524661135-423995f22d0b?w=1920&h=1080&fit=crop',

            # Yeni görseller - Modern Araçlar
            'ai_why': 'https://images.unsplash.com/photo-1677442136019-21780ecad995?w=1920&h=1080&fit=crop',
            'use_cases': 'https://images.unsplash.com/photo-1522542550221-31fd19575a2d?w=1920&h=1080&fit=crop',
            'workflow': 'https://images.unsplash.com/photo-1507238691740-187a5b1d37b8?w=1920&h=1080&fit=crop',
            'demo_setup': 'https://images.unsplash.com/photo-1573164713714-d95e436ab8d6?w=1920&h=1080&fit=crop',
            'best_practices': 'https://images.unsplash.com/photo-1553028826-f4804a6dba3b?w=1920&h=1080&fit=crop',
            'compare_tools': 'https://images.unsplash.com/photo-1551434678-e076c223a692?w=1920&h=1080&fit=crop',

            # Yeni görseller - Fonlama
            'funding_why': 'https://images.unsplash.com/photo-1579621970563-ebec7560ff3e?w=1920&h=1080&fit=crop',
            'tubitak_detail': 'https://images.unsplash.com/photo-1454165804606-c3d57bc86b40?w=1920&h=1080&fit=crop',
            'kosgeb_detail': 'https://images.unsplash.com/photo-1559526324-593bc073d938?w=1920&h=1080&fit=crop',
            'success_stories': 'https://images.unsplash.com/photo-1556761175-4b46a572b786?w=1920&h=1080&fit=crop',
            'calendar': 'https://images.unsplash.com/photo-1506784365847-bbad939e9335?w=1920&h=1080&fit=crop',

            # Yeni görseller - Mola
            'account_setup': 'https://images.unsplash.com/photo-1531482615713-2afd69097998?w=1920&h=1080&fit=crop',

            # Yeni görseller - Atölye
            'roadmap': 'https://images.unsplash.com/photo-1552664730-d307ca884978?w=1920&h=1080&fit=crop',
            'team_building': 'https://images.unsplash.com/photo-1556761175-5973dc0f32e7?w=1920&h=1080&fit=crop',
            'problem_criteria': 'https://images.unsplash.com/photo-1516321318423-f06f85e504b3?w=1920&h=1080&fit=crop',
            'swot_analysis': 'https://images.unsplash.com/photo-1454165804606-c3d57bc86b40?w=1920&h=1080&fit=crop',
            'template_demo': 'https://images.unsplash.com/photo-1586281380349-632531db7ed4?w=1920&h=1080&fit=crop',
            'research_methods': 'https://images.unsplash.com/photo-1532619675605-1ede6c2ed2b0?w=1920&h=1080&fit=crop',
            'problem_examples': 'https://images.unsplash.com/photo-1517245386807-bb43f82c33c4?w=1920&h=1080&fit=crop',
            'solution_matrix': 'https://images.unsplash.com/photo-1542744173-8e7e53415bb0?w=1920&h=1080&fit=crop',
            'mentor_tips': 'https://images.unsplash.com/photo-1573164713714-d95e436ab8d6?w=1920&h=1080&fit=crop',
            'pitch_template': 'https://images.unsplash.com/photo-1475721027785-f74eccf877e2?w=1920&h=1080&fit=crop',
            'feedback_loop': 'https://images.unsplash.com/photo-1556155092-490a1ba16284?w=1920&h=1080&fit=crop',

            # Yeni görseller - Kapanış
            'next_steps': 'https://images.unsplash.com/photo-1484788984921-03950022c9ef?w=1920&h=1080&fit=crop',
            'resources': 'https://images.unsplash.com/photo-1481627834876-b7833e8f5570?w=1920&h=1080&fit=crop',
        }

    def analyze_image_brightness(self, img):
        """Görselin parlaklığını analiz et ve optimal text rengi/overlay belirle"""
        # Görseli küçült (hız için)
        img_small = img.resize((200, 150))

        # RGB'ye çevir
        if img_small.mode != 'RGB':
            img_small = img_small.convert('RGB')

        # Numpy array'e çevir
        pixels = np.array(img_small)

        # Ortalama parlaklık (0-255)
        avg_brightness = np.mean(pixels)

        # RGB kanalları ortalaması
        avg_r = np.mean(pixels[:, :, 0])
        avg_g = np.mean(pixels[:, :, 1])
        avg_b = np.mean(pixels[:, :, 2])

        # Sonuçları döndür
        return {
            'brightness': avg_brightness,
            'is_dark': avg_brightness < 100,  # Karanlık görsel
            'is_light': avg_brightness > 155,  # Açık görsel
            'avg_r': avg_r,
            'avg_g': avg_g,
            'avg_b': avg_b
        }

    def get_optimal_text_style(self, img_analysis):
        """Tüm görselleri %50 karart"""
        # Tüm görseller için sabit %50 karartma
        return {
            'text_color': RGBColor(255, 255, 255),
            'overlay_needed': True,
            'overlay_opacity': 0.5,  # %50 karartma = brightness 0.5
            'shadow': False
        }

    def add_bg_image(self, slide, img_key, overlay=0.5):
        """Tam ekran arka plan görseli ekle"""
        self.stats['total'] += 1
        try:
            print(f"📥 İndiriliyor: {self.images[img_key][:50]}...")

            headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
            r = requests.get(self.images[img_key], timeout=15, headers=headers)
            r.raise_for_status()

            img = Image.open(BytesIO(r.content))

            # RGB'ye çevir
            if img.mode != 'RGB':
                if img.mode in ('RGBA', 'LA', 'P'):
                    bg = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode in ('RGBA', 'LA'):
                        bg.paste(img, mask=img.split()[-1])
                    else:
                        bg.paste(img)
                    img = bg
                else:
                    img = img.convert('RGB')

            # Görseli analiz et
            img_analysis = self.analyze_image_brightness(img)
            text_style = self.get_optimal_text_style(img_analysis)

            # Blur + kontrast
            img = img.filter(ImageFilter.GaussianBlur(radius=1))
            enhancer = ImageEnhance.Contrast(img)
            img = enhancer.enhance(1.2)

            # Akıllı Karartma - overlay yerine görseli direkt karart
            darken_info = "karartma yok"
            if text_style['overlay_needed'] and text_style['overlay_opacity'] > 0:
                # Görseli karart
                # overlay_opacity 0.6 = %40 karartma istiyoruz → brightness 0.6 kullan
                # overlay_opacity 0.75 = %25 karartma istiyoruz → brightness 0.75 kullan
                brightness_factor = text_style['overlay_opacity']
                brightness_enhancer = ImageEnhance.Brightness(img)
                img = brightness_enhancer.enhance(brightness_factor)
                darken_percent = int((1.0 - brightness_factor) * 100)
                darken_info = f"%{darken_percent} karartıldı"

            buf = BytesIO()
            img.save(buf, format='JPEG', quality=90)
            buf.seek(0)

            # Tam ekran görsel ekle
            left = top = Inches(0)
            width = self.prs.slide_width
            height = self.prs.slide_height

            pic = slide.shapes.add_picture(buf, left, top, width, height)

            # Görseli en arkaya al (z-order fix)
            sp_tree = slide.shapes._spTree
            sp_tree.remove(pic._element)
            sp_tree.insert(2, pic._element)  # 2. pozisyon (background'un hemen üstü)

            # OVERLAY ARTIK EKLENMİYOR - Görsel zaten karartıldı

            brightness = img_analysis['brightness']
            print(f"  ✅ Eklendi (parlaklık: {brightness:.0f}, {darken_info})\n")

            self.stats['success'] += 1

            # Progress bar ekle (eğer slide_counter varsa)
            if hasattr(self, 'slide_counter') and hasattr(self, 'total_slides'):
                self.add_progress_bar(slide, self.slide_counter, self.total_slides)

            # Text style'ı döndür (slayt oluştururken kullanılacak)
            return text_style
        except Exception as e:
            print(f"  ❌ HATA: {str(e)}\n")
            self.stats['failed'] += 1
            return False

    def add_title(self, slide, text, y=3.5, size=60, color='white'):
        """Başlık ekle"""
        box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(14), Inches(2))
        frame = box.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = self.colors[color] if color in self.colors else self.colors['white']

    def add_time_badge(self, slide, time_text, x=0.5, y=0.5):
        """Zaman rozeti"""
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(0.8))
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.colors['warning']
        badge.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.5), Inches(0.8))
        frame = text_box.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = frame.paragraphs[0]
        p.text = f"⏰ {time_text}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']

    def add_progress_bar(self, slide, current_slide, total_slides=61):
        """Progress bar ekle (alt kısımda)"""
        bar_width = 14  # inches
        bar_height = 0.15  # inches
        bar_x = 1  # inches
        bar_y = 8.7  # inches (alt kısımda)

        # Background bar (gri)
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bar_x), Inches(bar_y),
            Inches(bar_width), Inches(bar_height)
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = RGBColor(60, 60, 60)
        bg_bar.fill.transparency = 0.3
        bg_bar.line.fill.background()

        # Progress bar (renkli)
        progress_percentage = current_slide / total_slides
        progress_width = bar_width * progress_percentage

        if progress_width > 0:
            progress_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(bar_x), Inches(bar_y),
                Inches(progress_width), Inches(bar_height)
            )
            progress_bar.fill.solid()
            # Gradient renk: başlangıç mavi, orta mor, son yeşil
            if progress_percentage < 0.33:
                progress_bar.fill.fore_color.rgb = self.colors['primary']
            elif progress_percentage < 0.66:
                progress_bar.fill.fore_color.rgb = self.colors['purple']
            else:
                progress_bar.fill.fore_color.rgb = self.colors['success']
            progress_bar.line.fill.background()

        # Slide numarası (sağ alt köşe)
        slide_num_txt = slide.shapes.add_textbox(Inches(14.2), Inches(8.5), Inches(1.5), Inches(0.3))
        slide_num_txt.text_frame.paragraphs[0].text = f"{current_slide}/{total_slides}"
        slide_num_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        slide_num_txt.text_frame.paragraphs[0].font.size = Pt(14)
        slide_num_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

    def generate(self):
        """Sunumu oluştur"""
        print("\n" + "="*70)
        print("🎯 ULTİMATE 2 SAATLİK KAPSAMLI SUNUM OLUŞTURULUYOR...")
        print("="*70 + "\n")

        self.slide_counter = 0  # Progress bar için slide sayacı
        self.total_slides = 62  # Toplam slide sayısı (Design Thinking slaytı eklendi)

        # ===== AÇILIŞ (10:00-10:10) =====
        print("📍 BÖLÜM 1: Açılış...")

        # 1. Ana başlık
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'hero', 0.3)
        self.add_title(s, "SANAYİDE DİJİTAL DÖNÜŞÜM", 2.8, 76)
        sub = s.shapes.add_textbox(Inches(2), Inches(5), Inches(12), Inches(1))
        sub.text_frame.paragraphs[0].text = "AI Araçlarıyla Fikirden Ürüne"
        sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        sub.text_frame.paragraphs[0].font.size = Pt(32)
        sub.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

        # 2. Konuşmacı
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'team', 0.35)
        self.add_title(s, "Dr. Öğr. Üyesi Uğur CORUH", 2.5, 48)
        info = s.shapes.add_textbox(Inches(4), Inches(4), Inches(8), Inches(3))
        lines = ["RTEÜ Erasmus+", "19 Ekim 2025, Pazar", "10:00 - 12:00", "50 Öğrenci → 25 Proje"]
        for i, line in enumerate(lines):
            if i > 0:
                p = info.text_frame.add_paragraph()
            else:
                p = info.text_frame.paragraphs[0]
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(22)
            p.font.color.rgb = self.colors['white']

        # 3. Program akışı
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'technology', 0.3)
        self.add_title(s, "📅 BUGÜNÜN PROGRAMI", 0.8, 44)

        timeline = [
            {'time': '10:00-10:10', 'what': 'Açılış & Tanışma', 'color': self.colors['purple']},
            {'time': '10:10-10:30', 'what': 'Sanayi 4.0 & Dijital Dönüşüm', 'color': self.colors['primary']},
            {'time': '10:30-10:55', 'what': 'Modern Araçlar + CANLI DEMO', 'color': self.colors['success']},
            {'time': '10:55-11:20', 'what': 'Fonlama Ekosistemi', 'color': self.colors['warning']},
            {'time': '11:20-11:30', 'what': 'MOLA + Hazırlık', 'color': self.colors['pink']},
            {'time': '11:30-12:00', 'what': 'HANDS-ON ATÖLYE', 'color': self.colors['secondary']},
        ]

        y = Inches(2.2)
        for item in timeline:
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), y, Inches(12), Inches(1))
            box.fill.solid()
            box.fill.fore_color.rgb = item['color']
            box.line.fill.background()

            txt = s.shapes.add_textbox(Inches(2.3), y + Inches(0.2), Inches(11.4), Inches(0.6))
            txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = txt.text_frame.paragraphs[0]
            p.text = f"{item['time']} | {item['what']}"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            y += Inches(1.1)

        # 4. Ne öğreneceksiniz? (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'learning', 0.35)
        self.add_time_badge(s, "10:03")
        self.add_title(s, "📚 BUGÜN NE ÖĞRENECEKSİNİZ?", 1, 44)

        learnings = [
            {'icon': '🏭', 'title': 'Sanayi 4.0 Temelleri', 'detail': '8 temel teknoloji: IoT, AI, Big Data, Cloud, Dijital İkiz, Robotik, 3D Print, AR/VR'},
            {'icon': '🤖', 'title': 'AI Araçları Kullanımı', 'detail': 'Claude, ChatGPT, Gemini ile proje geliştirme, DeepResearch ile literatür tarama'},
            {'icon': '💰', 'title': 'Fonlama Fırsatları', 'detail': 'TÜBİTAK 2209-A (7.5K TL), KOSGEB (600K TL), Horizon Europe ve başvuru formatları'},
            {'icon': '📝', 'title': 'Proje Yazım Teknikleri', 'detail': 'TÜBİTAK formatı, 9 bölümlü proposal şablonu, AI ile otomatik döküman üretimi'},
            {'icon': '🎯', 'title': 'Problem-Çözüm Metodolojisi', 'detail': 'Gerçek sanayi sorunlarını belirleme, uygulanabilir çözümler geliştirme'},
            {'icon': '🚀', 'title': 'Hızlı Prototipleme', 'detail': 'AI yardımıyla 30 dakikada fikir → proje önerisi → pitch hazırlama'},
            {'icon': '🎤', 'title': '1 Dakikalık Pitch', 'detail': 'Problem-Çözüm-Etki formatında etkili sunum yapma becerileri'},
            {'icon': '🤝', 'title': 'Takım Çalışması & Networking', 'detail': 'İkili takımlar, peer feedback (3-2-1 metodu), mentörlük deneyimi'},
        ]

        y_start = Inches(2.7)
        row_height = Inches(1.2)
        col_width = Inches(7)
        margin = Inches(0.5)

        for i, item in enumerate(learnings):
            col = i % 2
            row = i // 2
            x_pos = Inches(0.8) + col * (col_width + margin)
            y_pos = y_start + row * (row_height + Inches(0.1))

            # Kart
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, col_width, row_height)
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(25, 25, 25)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['purple']
            card.line.width = Pt(2)

            # Icon + Başlık
            header_box = s.shapes.add_textbox(x_pos + Inches(0.15), y_pos + Inches(0.1), col_width - Inches(0.3), Inches(0.35))
            p = header_box.text_frame.paragraphs[0]
            p.text = f"{item['icon']} {item['title']}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            # Detay
            detail_box = s.shapes.add_textbox(x_pos + Inches(0.15), y_pos + Inches(0.5), col_width - Inches(0.3), Inches(0.65))
            detail_box.text_frame.word_wrap = True
            p2 = detail_box.text_frame.paragraphs[0]
            p2.text = item['detail']
            p2.font.size = Pt(20)
            p2.font.color.rgb = self.colors['light']

        # 5. Bugünün hedefleri (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'target_goals', 0.35)
        self.add_time_badge(s, "10:05")
        self.add_title(s, "🎯 BUGÜNÜN HEDEFLERİ", 1, 44)

        goals = [
            {'num': '1', 'goal': 'Sanayi 4.0 Farkındalığı', 'kpi': '8 temel teknolojiyi tanımak ve örneklemek'},
            {'num': '2', 'goal': 'AI Araçları Yetkinliği', 'kpi': 'En az 3 farklı AI aracını aktif kullanabilmek'},
            {'num': '3', 'goal': 'Fonlama Bilgisi', 'kpi': 'Minimum 5 fon kaynağını ve başvuru süreçlerini bilmek'},
            {'num': '4', 'goal': 'Proje Önerisi Oluşturma', 'kpi': 'TÜBİTAK formatında 1 tam proje önerisi hazırlamak'},
            {'num': '5', 'goal': 'Ekip İşbirliği & Pitch', 'kpi': '1 dakikalık profesyonel pitch sunumu yapabilmek'},
        ]

        y = Inches(2.4)
        for item in goals:
            # Ana kart
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(1.15))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(30, 30, 30)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['success']
            card.line.width = Pt(3)

            # Numara badge
            num_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), y + Inches(0.3), Inches(0.5), Inches(0.5))
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = self.colors['success']
            num_circle.line.fill.background()

            num_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.3), Inches(0.5), Inches(0.5))
            num_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = num_txt.text_frame.paragraphs[0]
            p.text = item['num']
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']

            # Hedef başlık
            goal_box = s.shapes.add_textbox(Inches(2.8), y + Inches(0.2), Inches(5), Inches(0.35))
            p = goal_box.text_frame.paragraphs[0]
            p.text = item['goal']
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            # KPI
            kpi_box = s.shapes.add_textbox(Inches(2.8), y + Inches(0.57), Inches(11), Inches(0.4))
            kpi_box.text_frame.word_wrap = True
            p = kpi_box.text_frame.paragraphs[0]
            p.text = f"✓ {item['kpi']}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['teal']

            y += Inches(1.2)

        # 6. Beklentiler vs Gerçeklik (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'expectations', 0.35)
        self.add_time_badge(s, "10:07")
        self.add_title(s, "🤔 BEKLENTİLER vs GERÇEKLİK", 1, 40)

        comparisons = [
            {
                'myth': '❌ "AI her şeyi yapar, biz hiçbir şey öğrenmeyiz"',
                'reality': '✅ AI yardımcıdır - Siz yönlendirir, kontrol eder, öğrenirsiniz'
            },
            {
                'myth': '❌ "Proje yazmak aylar sürer, çok zordur"',
                'reality': '✅ AI ile 30-60 dakikada ilk taslak, 1 haftada tamamlanır'
            },
            {
                'myth': '❌ "Fonlar sadece büyük şirketler ve deneyimlilere"',
                'reality': '✅ TÜBİTAK 2209-A sadece öğrencilere, KOSGEB herkese açık'
            },
            {
                'myth': '❌ "Sanayi 4.0 sadece mühendislik içindir"',
                'reality': '✅ Tüm sektörlere uygulanabilir: Sağlık, tarım, eğitim, perakende...'
            },
            {
                'myth': '❌ "Bugün sadece teori, uygulama yok"',
                'reality': '✅ 30 dakika hands-on atölye, kendi projenizi oluşturacaksınız!'
            },
        ]

        y = Inches(3.0)
        for comp in comparisons:
            # Container
            container = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y, Inches(14), Inches(1))
            container.fill.solid()
            container.fill.fore_color.rgb = RGBColor(25, 25, 25)
            container.fill.transparency = 0.2
            container.line.fill.background()

            # Myth (sol)
            myth_box = s.shapes.add_textbox(Inches(1.3), y + Inches(0.15), Inches(6.5), Inches(0.75))
            myth_box.text_frame.word_wrap = True
            p = myth_box.text_frame.paragraphs[0]
            p.text = comp['myth']
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['danger']

            # Arrow
            arrow_box = s.shapes.add_textbox(Inches(7.8), y + Inches(0.3), Inches(0.4), Inches(0.4))
            p = arrow_box.text_frame.paragraphs[0]
            p.text = "→"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.colors['warning']

            # Reality (sağ)
            reality_box = s.shapes.add_textbox(Inches(8.2), y + Inches(0.15), Inches(6.5), Inches(0.75))
            reality_box.text_frame.word_wrap = True
            p = reality_box.text_frame.paragraphs[0]
            p.text = comp['reality']
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['success']

            y += Inches(1.05)

        # ===== SANAYİ 4.0 (10:10-10:30) =====
        print("📍 BÖLÜM 2: Sanayi 4.0...")

        # 4. Problem
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'factory', 0.35)
        self.add_time_badge(s, "10:10-10:30")
        self.add_title(s, "SANAYİDE DİJİTAL DÖNÜŞÜM NEDEN?", 1.8, 44)

        stats = [
            {'n': '92%', 'l': 'Şirketler dönüşüm\nplanıyor'},
            {'n': '37%', 'l': 'Başarıyla\nuyguluyor'},
            {'n': '3.5X', 'l': 'Verimlilik\npotansiyeli'},
        ]

        x = Inches(2.5)
        for st in stats:
            c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4), Inches(3.5), Inches(3))
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(30, 30, 30)
            c.fill.transparency = 0.25
            c.line.color.rgb = self.colors['secondary']
            c.line.width = Pt(4)

            nb = s.shapes.add_textbox(x, Inches(4.5), Inches(3.5), Inches(1.5))
            nb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            nb.text_frame.paragraphs[0].text = st['n']
            nb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            nb.text_frame.paragraphs[0].font.size = Pt(64)
            nb.text_frame.paragraphs[0].font.bold = True
            nb.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

            lb = s.shapes.add_textbox(x, Inches(6), Inches(3.5), Inches(0.8))
            lb.text_frame.word_wrap = True
            lb.text_frame.paragraphs[0].text = st['l']
            lb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            lb.text_frame.paragraphs[0].font.size = Pt(20)
            lb.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            x += Inches(4)

        # 5. Sanayi 4.0 teknolojileri
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'network', 0.35)
        self.add_time_badge(s, "10:15")
        self.add_title(s, "🏭 SANAYİ 4.0: 8 TEKNOLOJI", 0.8, 44)

        techs = [
            {'i': '🔗', 'n': 'IoT', 'd': 'Akıllı sensörler, gerçek zamanlı veri'},
            {'i': '🧠', 'n': 'AI/ML', 'd': 'Tahmin, optimizasyon, otomasyon'},
            {'i': '📊', 'n': 'Big Data', 'd': 'Veri analitiği, pattern bulma'},
            {'i': '☁️', 'n': 'Cloud', 'd': 'Sınırsız kaynak, her yerden erişim'},
            {'i': '🔄', 'n': 'Dijital İkiz', 'd': 'Siber-fiziksel sistemler'},
            {'i': '🤖', 'n': 'Robotik', 'd': 'Otonom robotlar, cobots'},
            {'i': '🖨️', 'n': '3D Print', 'd': 'Hızlı prototipleme'},
            {'i': '🥽', 'n': 'AR/VR', 'd': 'Bakım, eğitim, montaj desteği'},
        ]

        y = Inches(2.7)
        row = 0
        for i, t in enumerate(techs):
            col = i % 2
            if i > 0 and col == 0:
                row += 1

            x = Inches(1.5) if col == 0 else Inches(8.5)
            yy = y + row * Inches(1.2)

            c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, Inches(6.5), Inches(1))
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(25, 25, 25)
            c.fill.transparency = 0.3
            c.line.color.rgb = self.colors['teal']
            c.line.width = Pt(2)

            txt = s.shapes.add_textbox(x + Inches(0.2), yy + Inches(0.15), Inches(6.1), Inches(0.7))
            txt.text_frame.word_wrap = True
            p = txt.text_frame.paragraphs[0]
            p.text = f"{t['i']} {t['n']}: {t['d']}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['white']

        # 7. Sanayi 4.0 Nedir? (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'industry_intro', 0.35)
        self.add_time_badge(s, "10:12")
        self.add_title(s, "🏭 SANAYİ 4.0 NEDİR?", 1, 44)

        intro_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2.5), Inches(12), Inches(1.5))
        intro_box.fill.solid()
        intro_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
        intro_box.fill.transparency = 0.2
        intro_box.line.color.rgb = self.colors['primary']
        intro_box.line.width = Pt(3)

        intro_txt = s.shapes.add_textbox(Inches(2.5), Inches(2.8), Inches(11), Inches(1))
        intro_txt.text_frame.word_wrap = True
        p = intro_txt.text_frame.paragraphs[0]
        p.text = "Üretim süreçlerinin dijitalleşmesi, otomasyon ve veri alışverişini içeren dördüncü sanayi devrimi. Fiziksel ve dijital sistemlerin birleşimi ile akıllı fabrikalar oluşturulur."
        p.font.size = Pt(20)
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER

        key_points = [
            {'icon': '🔗', 'title': 'Bağlantılılık', 'desc': 'Makineler, cihazlar ve sistemler birbirine bağlı'},
            {'icon': '📡', 'title': 'Gerçek Zamanlı Veri', 'desc': 'Anlık izleme ve analiz ile hızlı karar alma'},
            {'icon': '🤖', 'title': 'Otomasyon', 'desc': 'İnsan müdahalesi olmadan çalışan sistemler'},
            {'icon': '🧠', 'title': 'Yapay Zeka', 'desc': 'Öğrenen ve kendini geliştiren akıllı sistemler'},
        ]

        y = Inches(4.5)
        for item in key_points:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), y, Inches(12), Inches(0.7))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(25, 25, 25)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['teal']
            card.line.width = Pt(2)

            icon_txt = s.shapes.add_textbox(Inches(2.3), y + Inches(0.1), Inches(0.5), Inches(0.5))
            icon_txt.text_frame.paragraphs[0].text = item['icon']
            icon_txt.text_frame.paragraphs[0].font.size = Pt(24)

            title_txt = s.shapes.add_textbox(Inches(3), y + Inches(0.1), Inches(3), Inches(0.5))
            title_txt.text_frame.paragraphs[0].text = item['title']
            title_txt.text_frame.paragraphs[0].font.size = Pt(20)
            title_txt.text_frame.paragraphs[0].font.bold = True
            title_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            desc_txt = s.shapes.add_textbox(Inches(6.5), y + Inches(0.15), Inches(7), Inches(0.4))
            desc_txt.text_frame.word_wrap = True
            desc_txt.text_frame.paragraphs[0].text = item['desc']
            desc_txt.text_frame.paragraphs[0].font.size = Pt(20)
            desc_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            y += Inches(0.8)

        # 8. 4 Sanayi Devrimi (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'revolution', 0.35)
        self.add_time_badge(s, "10:14")
        self.add_title(s, "📜 4 SANAYİ DEVRİMİ", 0.8, 44)

        revolutions = [
            {'num': '1.0', 'year': '1784', 'tech': 'Buhar Makinesi', 'impact': 'Mekanik üretim, su ve buhar gücü', 'color': self.colors['secondary']},
            {'num': '2.0', 'year': '1870', 'tech': 'Elektrik', 'impact': 'Seri üretim, montaj hattı, elektrik enerjisi', 'color': self.colors['warning']},
            {'num': '3.0', 'year': '1969', 'tech': 'Bilgisayar & Otomasyon', 'impact': 'Dijital devrim, PLC, robotlar', 'color': self.colors['primary']},
            {'num': '4.0', 'year': '2011', 'tech': 'Siber-Fiziksel Sistemler', 'impact': 'IoT, AI, Big Data, Cloud, akıllı fabrikalar', 'color': self.colors['success']},
        ]

        y = Inches(2.7)
        for rev in revolutions:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(1.3))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(30, 30, 30)
            card.fill.transparency = 0.2
            card.line.color.rgb = rev['color']
            card.line.width = Pt(4)

            # Numara
            num_box = s.shapes.add_textbox(Inches(2), y + Inches(0.25), Inches(1.5), Inches(0.8))
            num_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = num_box.text_frame.paragraphs[0]
            p.text = rev['num']
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = rev['color']

            # Yıl
            year_box = s.shapes.add_textbox(Inches(3.8), y + Inches(0.1), Inches(1), Inches(0.4))
            p = year_box.text_frame.paragraphs[0]
            p.text = rev['year']
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['warning']

            # Teknoloji
            tech_box = s.shapes.add_textbox(Inches(3.8), y + Inches(0.5), Inches(4), Inches(0.6))
            tech_box.text_frame.word_wrap = True
            p = tech_box.text_frame.paragraphs[0]
            p.text = rev['tech']
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            # Etki
            impact_box = s.shapes.add_textbox(Inches(8.5), y + Inches(0.3), Inches(5.5), Inches(0.7))
            impact_box.text_frame.word_wrap = True
            p = impact_box.text_frame.paragraphs[0]
            p.text = rev['impact']
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['light']

            y += Inches(1.4)

        # 9-14. Her teknoloji için ayrı slayt (YENİ - 6 slayt)
        tech_details = [
            {
                'title': '🔗 IoT - Nesnelerin İnterneti',
                'time': '10:16',
                'img': 'iot_sensors',
                'definition': 'Fiziksel cihazların internet üzerinden veri toplaması ve paylaşması',
                'examples': [
                    '📡 Akıllı Sensörler: Sıcaklık, nem, titreşim, basınç sensörleri',
                    '🏭 Üretim Hattı İzleme: Makine performansı, enerji tüketimi',
                    '📦 Lojistik Takibi: RFID etiketleri, GPS izleme',
                    '⚡ Predictive Maintenance: Arıza öncesi uyarı sistemleri',
                ],
                'benefit': '💰 %25-30 bakım maliyeti düşüşü, %70 daha az beklenmedik duruş'
            },
            {
                'title': '🧠 AI/ML - Yapay Zeka & Makine Öğrenmesi',
                'time': '10:18',
                'img': 'artificial_intelligence',
                'definition': 'Sistemlerin veriden öğrenerek karar alması ve sürekli gelişmesi',
                'examples': [
                    '🔍 Kalite Kontrol: Görüntü işleme ile hatalı ürün tespiti',
                    '📈 Talep Tahmini: Satış verilerinden üretim planlaması',
                    '🤖 Süreç Optimizasyonu: En verimli üretim parametrelerini bulma',
                    '🎯 Enerji Yönetimi: Akıllı enerji tüketim optimizasyonu',
                ],
                'benefit': '🚀 %20-35 üretim verimliliği artışı, %99.9 kalite oranı'
            },
            {
                'title': '📊 Big Data - Büyük Veri Analitiği',
                'time': '10:20',
                'img': 'big_data',
                'definition': 'Büyük hacimli verinin toplanması, depolanması ve analiz edilmesi',
                'examples': [
                    '📉 Üretim Analitiği: Çevrim süresi, fire oranı, OEE metrikleri',
                    '🔬 Kalite Analitiği: Hata pattern analizi, root cause analysis',
                    '🌐 Tedarik Zinciri: Envanter optimizasyonu, risk analizi',
                    '👥 Müşteri Analitiği: Talep tahminleme, personalizasyon',
                ],
                'benefit': '📊 Karar alma hızında %5X artış, %15-20 maliyet düşüşü'
            },
            {
                'title': '☁️ Cloud Computing - Bulut Bilişim',
                'time': '10:22',
                'img': 'cloud_computing',
                'definition': 'İnternet üzerinden esnek ve ölçeklenebilir bilişim kaynakları',
                'examples': [
                    '💾 Veri Depolama: Sınırsız depolama, otomatik yedekleme',
                    '⚙️ Uygulama Hosting: ERP, MES, PLM sistemleri bulutta',
                    '🔐 Güvenli Erişim: Her yerden, her cihazdan güvenli bağlantı',
                    '📡 IoT Platform: Milyonlarca cihazdan veri toplama',
                ],
                'benefit': '💸 %30-40 IT altyapı maliyeti düşüşü, sonsuz ölçeklenebilirlik'
            },
            {
                'title': '🔄 Dijital İkiz (Digital Twin)',
                'time': '10:24',
                'img': 'digital_twin',
                'definition': 'Fiziksel varlıkların sanal ortamda birebir dijital kopyası',
                'examples': [
                    '🏭 Fabrika Simülasyonu: Üretim hattını sanal ortamda test',
                    '🔧 Ürün Geliştirme: Prototip öncesi sanal testler',
                    '📐 Süreç Optimizasyonu: "What-if" senaryoları simülasyonu',
                    '🎓 Eğitim: Risk almadan güvenli öğrenme ortamı',
                ],
                'benefit': '⚡ %50 daha hızlı ürün geliştirme, %75 test maliyeti düşüşü'
            },
            {
                'title': '🤖 Robotik & Otomasyon',
                'time': '10:26',
                'img': 'robotics',
                'definition': 'Otonom ve yarı-otonom sistemlerle üretim süreçlerinin otomasyonu',
                'examples': [
                    '🦾 Endüstriyel Robotlar: Kaynak, boyama, montaj robotları',
                    '🤝 Cobotlar (Collaborative): İnsanlarla yan yana çalışan robotlar',
                    '🚛 AGV/AMR: Otonom taşıma araçları, akıllı depolar',
                    '🎯 RPA: Ofis süreçlerinin otomasyonu (fatura, sipariş)',
                ],
                'benefit': '⏱️ %24/7 kesintisiz üretim, %60 daha yüksek hassasiyet'
            },
        ]

        for td in tech_details:
            s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            self.slide_counter += 1
            self.add_bg_image(s, td['img'], 0.35)
            self.add_time_badge(s, td['time'])
            self.add_title(s, td['title'], 0.8, 40)

            # Definition box - MOVED DOWN from 1.8" to 2.2"
            def_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(13), Inches(0.8))
            def_box.fill.solid()
            def_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
            def_box.fill.transparency = 0.2
            def_box.line.color.rgb = self.colors['primary']
            def_box.line.width = Pt(3)

            def_txt = s.shapes.add_textbox(Inches(2), Inches(2.4), Inches(12), Inches(0.5))
            def_txt.text_frame.word_wrap = True
            p = def_txt.text_frame.paragraphs[0]
            p.text = td['definition']
            p.font.size = Pt(20)
            p.font.italic = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER

            # Examples - MOVED DOWN from 3.0" to 3.2"
            y = Inches(3.2)
            for example in td['examples']:
                ex_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(0.75))
                ex_box.fill.solid()
                ex_box.fill.fore_color.rgb = RGBColor(25, 25, 25)
                ex_box.fill.transparency = 0.3
                ex_box.line.color.rgb = self.colors['teal']
                ex_box.line.width = Pt(2)

                ex_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.15), Inches(12), Inches(0.45))
                ex_txt.text_frame.word_wrap = True
                p = ex_txt.text_frame.paragraphs[0]
                p.text = example
                p.font.size = Pt(20)
                p.font.color.rgb = self.colors['white']

                y += Inches(0.75)  # REDUCED spacing from 0.85" to 0.75"

            # Benefit box - CHANGED from hardcoded Inches(7) to dynamic y position
            benefit_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y + Inches(0.1), Inches(13), Inches(0.9))
            benefit_box.fill.solid()
            benefit_box.fill.fore_color.rgb = self.colors['success']
            benefit_box.fill.transparency = 0.1
            benefit_box.line.color.rgb = self.colors['success']
            benefit_box.line.width = Pt(4)

            benefit_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.3), Inches(12), Inches(0.5))
            benefit_txt.text_frame.word_wrap = True
            p = benefit_txt.text_frame.paragraphs[0]
            p.text = f"✨ FAYDA: {td['benefit']}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER

        # 15. Türkiye'de Sanayi 4.0 (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'turkey_industry', 0.35)
        self.add_time_badge(s, "10:28")
        self.add_title(s, "🇹🇷 TÜRKİYE'DE SANAYİ 4.0", 1, 44)

        turkey_stats = [
            {'metric': 'Dijital Dönüşüm Hazırlığı', 'value': '52/100', 'rank': 'Avrupa 28. sıra'},
            {'metric': 'Sanayi 4.0 Yatırımı', 'value': '$2.5B', 'rank': '2020-2024 arası'},
            {'metric': 'Robot Yoğunluğu', 'value': '22/10K', 'rank': 'Çalışan başına (Dünya ort: 85)'},
            {'metric': 'Hedef 2030', 'value': 'Top 10', 'rank': 'Akıllı üretim endeksi'},
        ]

        y = Inches(3.0)
        for stat in turkey_stats:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), y, Inches(12), Inches(1.1))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(30, 30, 30)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['danger']
            card.line.width = Pt(3)

            metric_txt = s.shapes.add_textbox(Inches(2.5), y + Inches(0.15), Inches(5), Inches(0.4))
            p = metric_txt.text_frame.paragraphs[0]
            p.text = stat['metric']
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            value_txt = s.shapes.add_textbox(Inches(8), y + Inches(0.15), Inches(2.5), Inches(0.4))
            p = value_txt.text_frame.paragraphs[0]
            p.text = stat['value']
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.colors['secondary']
            p.alignment = PP_ALIGN.CENTER

            rank_txt = s.shapes.add_textbox(Inches(2.5), y + Inches(0.6), Inches(9), Inches(0.35))
            rank_txt.text_frame.word_wrap = True
            p = rank_txt.text_frame.paragraphs[0]
            p.text = f"→ {stat['rank']}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['teal']

            y += Inches(1.2)

        # Türkiye programları - MOVED DOWN from 7.3" to 7.9" to avoid overlap with cards
        programs_txt = s.shapes.add_textbox(Inches(2), Inches(7.9), Inches(12), Inches(1.0))
        programs_txt.text_frame.word_wrap = True
        p = programs_txt.text_frame.paragraphs[0]
        p.text = "🎯 PROGRAMLAR: Sanayi 4.0 Dönüşüm Programı, Dijital Türkiye 2023, Milli Teknoloji Hamlesi, KOSGEB Dijital Dönüşüm Desteği"
        p.font.size = Pt(20)
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER

        # ===== MODERN ARAÇLAR (10:30-10:55) =====
        print("📍 BÖLÜM 3: AI Araçları...")

        # 6. Araçlar girişi
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'code', 0.3)
        self.add_time_badge(s, "10:30-10:55")
        self.add_title(s, "🛠️ MODERN PROJE ARAÇLARI", 2.5, 56)
        sub = s.shapes.add_textbox(Inches(2), Inches(5), Inches(12), Inches(0.8))
        sub.text_frame.paragraphs[0].text = "Cursor • Claude Code • Claude • Gemini • ChatGPT • PlantUML • DrawIO • Python"
        sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        sub.text_frame.paragraphs[0].font.size = Pt(20)
        sub.text_frame.paragraphs[0].font.color.rgb = self.colors['success']

        # 7-10. Her araç kategorisi detaylı
        tools_detail = [
            {
                'title': '🤖 CLAUDE EKOSİSTEMİ',
                'time': '10:32',
                'img': 'ai_brain',
                'tools': [
                    {'name': 'Claude (chat)', 'use': 'Brainstorm, araştırma, problem çözme'},
                    {'name': 'Claude Projects', 'use': 'Custom instructions, döküman yükleme, proje yönetimi'},
                    {'name': 'Claude Code', 'use': 'Kod yazma, debugging, refactoring'},
                    {'name': 'DeepResearch', 'use': 'Kapsamlı literatür taraması, patent araştırması'},
                ]
            },
            {
                'title': '💻 GELİŞTİRME ARAÇLARI',
                'time': '10:40',
                'img': 'laptop',
                'tools': [
                    {'name': 'Cursor', 'use': 'AI pair programming, otomatik kod tamamlama'},
                    {'name': 'GitHub Copilot', 'use': 'Kod önerileri, test yazımı'},
                    {'name': 'GitHub', 'use': 'Versiyon kontrol, CI/CD, collaboration'},
                    {'name': 'VS Code', 'use': 'Code editor, extensions, terminal'},
                ]
            },
            {
                'title': '🎨 DİYAGRAM & DÖKÜMAN',
                'time': '10:48',
                'img': 'workspace',
                'tools': [
                    {'name': 'PlantUML', 'use': 'UML diyagramları, kod ile diagram'},
                    {'name': 'DrawIO', 'use': 'Akış diyagramları, mimari şemaları'},
                    {'name': 'Python-docx', 'use': 'Otomatik Word dökümanı'},
                    {'name': 'Python-pptx', 'use': 'Otomatik PowerPoint sunumu'},
                ]
            },
            {
                'title': '🔍 ARAŞTIRMA & YARDIMCI',
                'time': '10:52',
                'img': 'data',
                'tools': [
                    {'name': 'ChatGPT', 'use': 'Alternatif bakış açısı, yaratıcı fikirler'},
                    {'name': 'Gemini', 'use': 'Google entegrasyonu, multimodal analiz'},
                    {'name': 'Perplexity', 'use': 'Güncel bilgi, kaynak araştırması'},
                    {'name': 'Python', 'use': 'Otomasyon, veri işleme, web scraping'},
                ]
            }
        ]

        for td in tools_detail:
            s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            self.slide_counter += 1
            self.add_bg_image(s, td['img'], 0.35)
            self.add_time_badge(s, td['time'])
            self.add_title(s, td['title'], 1, 40)

            y = Inches(3.0)
            for tool in td['tools']:
                c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), y, Inches(12), Inches(1))
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(30, 30, 30)
                c.fill.transparency = 0.3
                c.line.color.rgb = self.colors['primary']
                c.line.width = Pt(2)

                nm = s.shapes.add_textbox(Inches(2.3), y + Inches(0.15), Inches(3), Inches(0.3))
                nm.text_frame.paragraphs[0].text = tool['name']
                nm.text_frame.paragraphs[0].font.size = Pt(20)
                nm.text_frame.paragraphs[0].font.bold = True
                nm.text_frame.paragraphs[0].font.color.rgb = self.colors['success']

                us = s.shapes.add_textbox(Inches(5.5), y + Inches(0.15), Inches(8.2), Inches(0.7))
                us.text_frame.word_wrap = True
                us.text_frame.paragraphs[0].text = f"→ {tool['use']}"
                us.text_frame.paragraphs[0].font.size = Pt(20)
                us.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

                y += Inches(1.15)

        # 11. CANLI DEMO
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'presentation', 0.5)
        self.add_time_badge(s, "10:53", x=6.5)
        self.add_title(s, "🎬 CANLI DEMO", 3, 72, 'secondary')
        sub = s.shapes.add_textbox(Inches(3), Inches(5.5), Inches(10), Inches(1))
        sub.text_frame.paragraphs[0].text = "Claude ile 5 dakikada proje önerisi"
        sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        sub.text_frame.paragraphs[0].font.size = Pt(32)
        sub.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

        # 12. AI Araçları Neden? (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'ai_why', 0.35)
        self.add_time_badge(s, "10:42")
        self.add_title(s, "🤔 AI ARAÇLARI NEDEN GEREKLİ?", 1, 42)

        why_items = [
            {'icon': '⚡', 'reason': 'Hız', 'detail': '100X daha hızlı kod yazma, 10 dakikada proje önerisi'},
            {'icon': '🎯', 'reason': 'Doğruluk', 'detail': 'Syntax hataları yok, best practice kodlama'},
            {'icon': '📚', 'reason': 'Öğrenme', 'detail': 'Yeni diller/framework\'leri anında öğrenme'},
            {'icon': '💡', 'reason': 'Yaratıcılık', 'detail': 'Farklı yaklaşımlar, alternatif çözümler'},
            {'icon': '🔄', 'reason': '24/7 Asistan', 'detail': 'Gece 3\'te bile yardım alabilen AI pair programmer'},
            {'icon': '🌍', 'reason': 'Demokratizasyon', 'detail': 'Herkes kod yazabilir, herkes proje geliştirebilir'},
        ]

        y = Inches(2.5)  # CHANGED from 3.0" to 2.5" - kutuları yukarı al (AI ARAÇLARI NEDEN GEREKLİ)
        for item in why_items:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(0.9))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(30, 30, 30)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['success']
            card.line.width = Pt(3)

            icon_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.15), Inches(0.6), Inches(0.6))
            icon_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = icon_txt.text_frame.paragraphs[0]
            p.text = item['icon']
            p.font.size = Pt(32)
            p.alignment = PP_ALIGN.CENTER

            reason_txt = s.shapes.add_textbox(Inches(2.8), y + Inches(0.15), Inches(2.5), Inches(0.6))
            reason_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = reason_txt.text_frame.paragraphs[0]
            p.text = item['reason']
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            detail_txt = s.shapes.add_textbox(Inches(5.5), y + Inches(0.2), Inches(8.5), Inches(0.5))
            detail_txt.text_frame.word_wrap = True
            p = detail_txt.text_frame.paragraphs[0]
            p.text = item['detail']
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['light']

            y += Inches(1)

        # 13. AI Kullanım Senaryoları (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'use_cases', 0.35)
        self.add_time_badge(s, "10:44")
        self.add_title(s, "🎯 AI KULLANIM SENARYOLARI", 0.8, 42)

        scenarios = [
            {'task': 'Proje Fikri Bulma', 'before': '2-3 hafta beyin fırtınası', 'after': '30 dakika Claude ile', 'gain': '%95 zaman'},
            {'task': 'Literatür Taraması', 'before': '1-2 hafta manuel araştırma', 'after': '1 saat DeepResearch', 'gain': '%90 zaman'},
            {'task': 'Kod Yazımı', 'before': '40 sat/hafta manuel kodlama', 'after': '10 sat/hafta AI ile', 'gain': '%75 zaman'},
            {'task': 'Dokümantasyon', 'before': '1 hafta Word\'de yazma', 'after': '2 saat python-docx', 'gain': '%85 zaman'},
            {'task': 'Sunum Hazırlama', 'before': '2 gün PowerPoint', 'after': '1 saat python-pptx', 'gain': '%90 zaman'},
        ]

        y = Inches(2.7)
        for sc in scenarios:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(14.4), Inches(1))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(25, 25, 25)
            card.fill.transparency = 0.2
            card.line.color.rgb = self.colors['primary']
            card.line.width = Pt(2)

            task_txt = s.shapes.add_textbox(Inches(1.2), y + Inches(0.15), Inches(3), Inches(0.7))
            task_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = task_txt.text_frame.paragraphs[0]
            p.text = sc['task']
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            before_txt = s.shapes.add_textbox(Inches(4.5), y + Inches(0.25), Inches(3.5), Inches(0.5))
            before_txt.text_frame.word_wrap = True
            p = before_txt.text_frame.paragraphs[0]
            p.text = f"❌ {sc['before']}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['danger']

            arrow = s.shapes.add_textbox(Inches(8.2), y + Inches(0.3), Inches(0.5), Inches(0.4))
            p = arrow.text_frame.paragraphs[0]
            p.text = "→"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['warning']
            p.alignment = PP_ALIGN.CENTER

            after_txt = s.shapes.add_textbox(Inches(8.8), y + Inches(0.25), Inches(3.5), Inches(0.5))
            after_txt.text_frame.word_wrap = True
            p = after_txt.text_frame.paragraphs[0]
            p.text = f"✅ {sc['after']}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['success']

            gain_badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(12.5), y + Inches(0.25), Inches(2.2), Inches(0.5))
            gain_badge.fill.solid()
            gain_badge.fill.fore_color.rgb = self.colors['success']
            gain_badge.line.fill.background()

            gain_txt = s.shapes.add_textbox(Inches(12.5), y + Inches(0.25), Inches(2.2), Inches(0.5))
            gain_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = gain_txt.text_frame.paragraphs[0]
            p.text = f"💰 {sc['gain']}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']
            p.alignment = PP_ALIGN.CENTER

            y += Inches(1.1)

        # 14. AI İş Akışı (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'workflow', 0.35)
        self.add_time_badge(s, "10:46")
        self.add_title(s, "🔄 AI İLE PROJE GELİŞTİRME AKIM", 0.8, 40)

        workflow_steps = [
            {'num': '1', 'step': 'İdea', 'tool': 'Claude Chat', 'action': 'Beyin fırtınası, problem tanımlama'},
            {'num': '2', 'step': 'Research', 'tool': 'DeepResearch', 'action': 'Literatür tarama, patent araştırma'},
            {'num': '3', 'step': 'Plan', 'tool': 'Claude Projects', 'action': 'Proje outline, görev dağılımı'},
            {'num': '4', 'step': 'Code', 'tool': 'Cursor/Claude Code', 'action': 'Kod yazma, debugging, refactoring'},
            {'num': '5', 'step': 'Document', 'tool': 'Python-docx', 'action': 'Otomatik Word/PDF döküman'},
            {'num': '6', 'step': 'Present', 'tool': 'Python-pptx', 'action': 'Otomatik PowerPoint sunumu'},
            {'num': '7', 'step': 'Deploy', 'tool': 'GitHub + CI/CD', 'action': 'Versiyon kontrol, otomatik deploy'},
        ]

        y = Inches(2.3)  # CHANGED from 2.0" to 2.3" - kutuları az aşağı al (AI PROJE GELİŞTİRME AKIM)
        for wf in workflow_steps:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), y, Inches(12), Inches(0.75))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(30, 30, 30)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['teal']
            card.line.width = Pt(2)

            num_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.3), y + Inches(0.15), Inches(0.45), Inches(0.45))
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = self.colors['secondary']
            num_circle.line.fill.background()

            num_txt = s.shapes.add_textbox(Inches(2.3), y + Inches(0.15), Inches(0.45), Inches(0.45))
            num_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = num_txt.text_frame.paragraphs[0]
            p.text = wf['num']
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER

            step_txt = s.shapes.add_textbox(Inches(2.9), y + Inches(0.15), Inches(1.5), Inches(0.45))
            step_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = step_txt.text_frame.paragraphs[0]
            p.text = wf['step']
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            tool_txt = s.shapes.add_textbox(Inches(4.6), y + Inches(0.15), Inches(3), Inches(0.45))
            tool_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tool_txt.text_frame.paragraphs[0]
            p.text = f"🔧 {wf['tool']}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['success']

            action_txt = s.shapes.add_textbox(Inches(7.8), y + Inches(0.2), Inches(5.8), Inches(0.35))
            action_txt.text_frame.word_wrap = True
            p = action_txt.text_frame.paragraphs[0]
            p.text = wf['action']
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['light']

            y += Inches(0.8)

        # 15. Demo Hazırlığı (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'demo_setup', 0.3)
        self.add_time_badge(s, "10:48")
        self.add_title(s, "🎬 DEMO HAZIRLIĞI", 1, 48)

        demo_prep = [
            '✅ Claude hesabı açın (claude.ai - ücretsiz)',
            '✅ Örnek problem hazırlayın: "Fabrikada enerji israfı"',
            '✅ Prompt şablonu: "TÜBİTAK 2209-A formatında proje önerisi yaz"',
            '✅ İzleyicilere göstereceğiniz ekranı hazırlayın',
            '✅ 5 dakika süre tutun - hız önemli!',
        ]

        y = Inches(3)
        for i, prep in enumerate(demo_prep):
            prep_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), y, Inches(10), Inches(0.7))
            prep_box.fill.solid()
            prep_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
            prep_box.fill.transparency = 0.2
            prep_box.line.color.rgb = self.colors['success']
            prep_box.line.width = Pt(3)

            prep_txt = s.shapes.add_textbox(Inches(3.3), y + Inches(0.15), Inches(9.4), Inches(0.4))
            prep_txt.text_frame.word_wrap = True
            p = prep_txt.text_frame.paragraphs[0]
            p.text = prep
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            y += Inches(0.8)

        # 16. En İyi Uygulamalar (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'best_practices', 0.35)
        self.add_time_badge(s, "10:50")
        self.add_title(s, "⭐ AI KULLANIM İPUÇLARI", 0.8, 42)

        tips = [
            {'do': '✅ Spesifik olun', 'dont': '❌ Belirsiz sorular', 'example': '"Python Flask API yaz" yerine "Flask REST API, JWT auth, SQLite DB"'},
            {'do': '✅ Bağlam verin', 'dont': '❌ Tek satır prompt', 'example': 'Projenin amacını, kullanıcı tipini, kısıtlamaları belirtin'},
            {'do': '✅ İteratif çalışın', 'dont': '❌ İlk çıktıyı kabul edin', 'example': '"Şunu ekle, bunu değiştir" diyerek geliştirin'},
            {'do': '✅ Doğrulayın', 'dont': '❌ Kör güven', 'example': 'AI\'ın verdiği kodu test edin, anlamadığınızı sorun'},
            {'do': '✅ Versiyon kontrolü', 'dont': '❌ Sadece kopyala-yapıştır', 'example': 'Git kullanın, her değişikliği commit edin'},
        ]

        y = Inches(2.7)
        for tip in tips:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(14.4), Inches(1.1))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(25, 25, 25)
            card.fill.transparency = 0.2
            card.line.color.rgb = self.colors['primary']
            card.line.width = Pt(2)

            do_txt = s.shapes.add_textbox(Inches(1.2), y + Inches(0.1), Inches(3), Inches(0.35))
            p = do_txt.text_frame.paragraphs[0]
            p.text = tip['do']
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['success']

            dont_txt = s.shapes.add_textbox(Inches(4.5), y + Inches(0.1), Inches(3.5), Inches(0.35))
            p = dont_txt.text_frame.paragraphs[0]
            p.text = tip['dont']
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['danger']

            example_txt = s.shapes.add_textbox(Inches(1.2), y + Inches(0.5), Inches(13.5), Inches(0.55))
            example_txt.text_frame.word_wrap = True
            p = example_txt.text_frame.paragraphs[0]
            p.text = f"💡 {tip['example']}"
            p.font.size = Pt(20)
            p.font.italic = True
            p.font.color.rgb = self.colors['teal']

            y += Inches(1.2)

        # 17. Araç Karşılaştırması (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'compare_tools', 0.35)
        self.add_time_badge(s, "10:52")
        self.add_title(s, "⚖️ AI ARAÇLARI KARŞILAŞTIRMA", 0.8, 38)

        comparisons = [
            {'tool': 'Claude', 'strength': 'En uzun context (200K), Projects', 'weakness': 'Kod çalıştıramaz', 'best_for': 'Araştırma, yazma'},
            {'tool': 'ChatGPT', 'strength': 'Plugin ekosistemi, DALL-E entegre', 'weakness': 'Kısa context (128K)', 'best_for': 'Genel amaçlı'},
            {'tool': 'Gemini', 'strength': 'Google entegrasyonu, ücretsiz', 'weakness': 'Kod kalitesi düşük', 'best_for': 'Hızlı sorular'},
            {'tool': 'Cursor', 'strength': 'IDE entegre, kod tamamlama', 'weakness': 'Ücretli ($20/ay)', 'best_for': 'Yazılım geliştirme'},
            {'tool': 'Claude Code', 'strength': 'Terminal komutları, dosya edit', 'weakness': 'Beta aşaması', 'best_for': 'Full-stack projeler'},
        ]

        y = Inches(2.7)
        for comp in comparisons:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y, Inches(14), Inches(1))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(30, 30, 30)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['purple']
            card.line.width = Pt(2)

            tool_txt = s.shapes.add_textbox(Inches(1.5), y + Inches(0.15), Inches(2), Inches(0.7))
            tool_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tool_txt.text_frame.paragraphs[0]
            p.text = comp['tool']
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            strength_txt = s.shapes.add_textbox(Inches(3.7), y + Inches(0.05), Inches(4.5), Inches(0.4))
            strength_txt.text_frame.word_wrap = True
            p = strength_txt.text_frame.paragraphs[0]
            p.text = f"✅ {comp['strength']}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['success']

            weakness_txt = s.shapes.add_textbox(Inches(3.7), y + Inches(0.48), Inches(4.5), Inches(0.4))
            weakness_txt.text_frame.word_wrap = True
            p = weakness_txt.text_frame.paragraphs[0]
            p.text = f"❌ {comp['weakness']}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['danger']

            best_badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), y + Inches(0.25), Inches(6), Inches(0.5))
            best_badge.fill.solid()
            best_badge.fill.fore_color.rgb = self.colors['teal']
            best_badge.fill.transparency = 0.2
            best_badge.line.color.rgb = self.colors['teal']
            best_badge.line.width = Pt(2)

            best_txt = s.shapes.add_textbox(Inches(8.7), y + Inches(0.3), Inches(5.6), Inches(0.4))
            best_txt.text_frame.word_wrap = True
            p = best_txt.text_frame.paragraphs[0]
            p.text = f"🎯 En iyi: {comp['best_for']}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            y += Inches(1.05)

        # ===== FONLAMA (10:55-11:20) =====
        print("📍 BÖLÜM 4: Fonlama...")

        # 12. Fonlama girişi
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'money', 0.3)
        self.add_time_badge(s, "10:55-11:20")
        self.add_title(s, "💰 FONLAMA EKOSİSTEMİ", 3, 60)

        # 13. Türkiye fonları
        tr_funds = [
            {'n': 'TÜBİTAK 2209-A', 'a': '7.500 TL', 't': '12 ay', 'i': 'Öğrenci projeleri, mentor gerekli'},
            {'n': 'KOSGEB Ar-Ge', 'a': '600K TL', 't': '24 ay', 'i': 'Girişimler, %75 hibe, prototip'},
            {'n': 'Teknofest', 'a': '100K TL', 't': '8 ay', 'i': '40+ kategori, prototip zorunlu'},
            {'n': 'TÜBİTAK 1512', 'a': '750K TL', 't': '18 ay', 'i': 'Teknogirişim, 500K hibe'},
        ]

        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'growth', 0.35)
        self.add_time_badge(s, "10:58")
        self.add_title(s, "🇹🇷 TÜRKİYE FONLARI", 1, 40)

        y = Inches(3.0)
        for f in tr_funds:
            c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(1.3))
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(30, 30, 30)
            c.fill.transparency = 0.25
            c.line.color.rgb = self.colors['warning']
            c.line.width = Pt(3)

            nm = s.shapes.add_textbox(Inches(2), y + Inches(0.15), Inches(3.5), Inches(0.4))
            nm.text_frame.paragraphs[0].text = f['n']
            nm.text_frame.paragraphs[0].font.size = Pt(20)
            nm.text_frame.paragraphs[0].font.bold = True
            nm.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            meta = s.shapes.add_textbox(Inches(6), y + Inches(0.15), Inches(3), Inches(0.4))
            meta.text_frame.paragraphs[0].text = f"💵 {f['a']} • ⏱️ {f['t']}"
            meta.text_frame.paragraphs[0].font.size = Pt(20)
            meta.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

            info = s.shapes.add_textbox(Inches(2), y + Inches(0.65), Inches(11.5), Inches(0.5))
            info.text_frame.word_wrap = True
            info.text_frame.paragraphs[0].text = f['i']
            info.text_frame.paragraphs[0].font.size = Pt(20)
            info.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            y += Inches(1.45)

        # 14. Uluslararası
        intl = [
            {'n': 'EU Horizon', 'a': '1-10M EUR', 'i': 'Konsorsiyum gerekli, TRL 3-9'},
            {'n': 'Erasmus+ KA2', 'a': '60-450K EUR', 'i': 'Eğitim ortaklıkları, 2-3 yıl'},
            {'n': 'EIC Accelerator', 'a': '2.5M EUR', 'i': 'Deep-tech startups, equity'},
        ]

        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'future', 0.35)
        self.add_time_badge(s, "11:08")
        self.add_title(s, "🌍 ULUSLARARASI FONLAR", 1, 40)

        y = Inches(2.3)  # CHANGED from 2.5" to 2.3" - daha da yukarı (SLIDE 24)
        for f in intl:
            c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), y, Inches(12), Inches(1.5))
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(30, 30, 30)
            c.fill.transparency = 0.25
            c.line.color.rgb = self.colors['primary']
            c.line.width = Pt(3)

            nm = s.shapes.add_textbox(Inches(2.5), y + Inches(0.2), Inches(4), Inches(0.5))
            nm.text_frame.paragraphs[0].text = f['n']
            nm.text_frame.paragraphs[0].font.size = Pt(22)
            nm.text_frame.paragraphs[0].font.bold = True
            nm.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            amt = s.shapes.add_textbox(Inches(7), y + Inches(0.2), Inches(3), Inches(0.5))
            amt.text_frame.paragraphs[0].text = f"💶 {f['a']}"
            amt.text_frame.paragraphs[0].font.size = Pt(20)
            amt.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

            info = s.shapes.add_textbox(Inches(2.5), y + Inches(0.85), Inches(9), Inches(0.5))
            info.text_frame.word_wrap = True
            info.text_frame.paragraphs[0].text = f['i']
            info.text_frame.paragraphs[0].font.size = Pt(20)
            info.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            y += Inches(1.7)

        # 15. Fonlama neden önemli? (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'funding_why', 0.35)
        self.add_time_badge(s, "11:10")
        self.add_title(s, "🎯 FONLAMA NEDEN ÖNEMLİ?", 1, 44)

        why_funding = [
            {'icon': '💰', 'title': 'Kaynak Sağlar', 'detail': 'Prototip, ekipman, yazılım lisansları, danışmanlık hizmetleri'},
            {'icon': '🚀', 'title': 'Hızlandırır', 'detail': 'Fikir aşamasından ürüne 2-3 yıl yerine 6-12 ayda'},
            {'icon': '✅', 'title': 'Meşruiyet', 'detail': 'TÜBİTAK/KOSGEB onayı, yatırımcı güveni, referans olur'},
            {'icon': '🎓', 'title': 'Öğrenme', 'detail': 'Proje yönetimi, bütçe planlama, raporlama disiplini'},
            {'icon': '🌐', 'title': 'Ağ Genişletir', 'detail': 'Mentorlar, sektör uzmanları, diğer girişimcilerle tanışma'},
            {'icon': '📈', 'title': 'Değer Kanıtı', 'detail': 'MVP çıktısı, pilot müşteri, pazar validasyonu'},
        ]

        cols = 2
        x_start = 1.5
        y_start = 2.8
        w = 5.8
        h = 1.6
        x_gap = 0.4
        y_gap = 0.3

        for i, item in enumerate(why_funding):
            row = i // cols
            col = i % cols
            x = x_start + col * (w + x_gap)
            y = y_start + row * (h + y_gap)

            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(25, 25, 25)
            card.fill.transparency = 0.2
            card.line.color.rgb = self.colors['success']
            card.line.width = Pt(2)

            icon_box = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(0.8), Inches(0.8))
            icon_box.text_frame.paragraphs[0].text = item['icon']
            icon_box.text_frame.paragraphs[0].font.size = Pt(36)

            title_box = s.shapes.add_textbox(Inches(x + 1.1), Inches(y + 0.2), Inches(w - 1.3), Inches(0.5))
            title_box.text_frame.paragraphs[0].text = item['title']
            title_box.text_frame.paragraphs[0].font.size = Pt(20)
            title_box.text_frame.paragraphs[0].font.bold = True
            title_box.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            detail_box = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.8), Inches(w - 0.6), Inches(0.7))
            detail_box.text_frame.word_wrap = True
            detail_box.text_frame.paragraphs[0].text = item['detail']
            detail_box.text_frame.paragraphs[0].font.size = Pt(20)
            detail_box.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

        # 16. TÜBİTAK 2209-A Detaylı (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'tubitak_detail', 0.35)
        self.add_time_badge(s, "11:12")
        self.add_title(s, "📋 TÜBİTAK 2209-A BAŞVURU SÜRECİ", 1, 42)

        tubitak_steps = [
            {'num': '1', 'title': 'Ön Hazırlık (Ocak)', 'tasks': ['Proje ekibi oluştur (2-3 kişi)', 'Akademik danışman bul', 'Fikir netleştir'], 'duration': '2 hafta'},
            {'num': '2', 'title': 'Başvuru (Şubat-Mart)', 'tasks': ['Proje özeti (500 kelime)', 'Detaylı plan (10 sayfa)', 'Bütçe tablosu', 'Özgeçmişler'], 'duration': '3 hafta'},
            {'num': '3', 'title': 'Değerlendirme (Nisan)', 'tasks': ['Bilimsel değerlendirme', 'Bütçe inceleme', 'Panel görüşmesi (seçilenlerde)'], 'duration': '4 hafta'},
            {'num': '4', 'title': 'Sonuç & Sözleşme (Mayıs)', 'tasks': ['Sonuç açıklama', 'Sözleşme imzalama', 'Proje başlangıcı'], 'duration': '2 hafta'},
            {'num': '5', 'title': 'Proje Yürütme (12 ay)', 'tasks': ['6 aylık ara rapor', 'Harcama belgeleri', 'Nihai rapor & sunum'], 'duration': '12 ay'},
        ]

        y = Inches(2.5)  # CHANGED from 3.1" to 2.5" - adjusted position (SLIDE 26)
        for step in tubitak_steps:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(1.15))  # CHANGED height 1.25"->1.15" - daha basık
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['primary']
            card.line.width = Pt(3)

            # Numara - adjusted for smaller card
            num_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.8), y + Inches(0.28), Inches(0.6), Inches(0.6))
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = self.colors['primary']
            num_circle.line.fill.background()

            num_txt = s.shapes.add_textbox(Inches(1.8), y + Inches(0.28), Inches(0.6), Inches(0.6))
            num_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            num_txt.text_frame.paragraphs[0].text = step['num']
            num_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_txt.text_frame.paragraphs[0].font.size = Pt(24)
            num_txt.text_frame.paragraphs[0].font.bold = True
            num_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Başlık
            title_txt = s.shapes.add_textbox(Inches(2.6), y + Inches(0.15), Inches(5), Inches(0.4))
            title_txt.text_frame.paragraphs[0].text = step['title']
            title_txt.text_frame.paragraphs[0].font.size = Pt(20)
            title_txt.text_frame.paragraphs[0].font.bold = True
            title_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Süre
            dur_txt = s.shapes.add_textbox(Inches(12.5), y + Inches(0.4), Inches(1.5), Inches(0.4))
            dur_txt.text_frame.paragraphs[0].text = f"⏱️ {step['duration']}"
            dur_txt.text_frame.paragraphs[0].font.size = Pt(20)
            dur_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

            # Görevler
            tasks_txt = s.shapes.add_textbox(Inches(2.6), y + Inches(0.6), Inches(9.5), Inches(0.6))
            tasks_txt.text_frame.word_wrap = True
            tasks_txt.text_frame.paragraphs[0].text = " • ".join(step['tasks'])
            tasks_txt.text_frame.paragraphs[0].font.size = Pt(20)
            tasks_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            y += Inches(1.25)  # CHANGED from 1.35" to 1.25" - daha kompakt

        # 17. KOSGEB Detaylı (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'kosgeb_detail', 0.35)
        self.add_time_badge(s, "11:14")
        self.add_title(s, "📋 KOSGEB AR-GE BAŞVURU SÜRECİ", 1, 42)

        kosgeb_steps = [
            {'phase': 'Ön Koşullar', 'items': ['18-29 yaş arası + Mezun olma', 'Şirket kuruluşu (Ltd/AŞ)', 'KOSGEB Girişimcilik eğitimi'], 'color': self.colors['warning']},
            {'phase': 'Başvuru Dönemi (Yılda 2 kez)', 'items': ['İş planı hazırlama (20-30 sayfa)', 'Teknik dökümantasyon', 'Finansal tablolar', 'Online başvuru'], 'color': self.colors['primary']},
            {'phase': 'Değerlendirme (3 ay)', 'items': ['Ön değerlendirme (40 puan)', 'Yazılı sınav (online)', 'Mülakat (30 puan)', 'Sıralama ve sonuç'], 'color': self.colors['secondary']},
            {'phase': 'Destek Aşamaları', 'items': ['İşletme Kuruluş: 50K TL', 'Makine/Teçhizat: 450K TL', 'İşletme Sermayesi: 100K TL'], 'color': self.colors['success']},
        ]

        y = Inches(2.2)  # CHANGED to 2.2" - başlangıç pozisyonu optimize
        for step in kosgeb_steps:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(1.5))  # CHANGED: x=1.5" (daha geniş), height=1.5" (daha kompakt)
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(25, 25, 25)
            card.fill.transparency = 0.2
            card.line.color.rgb = step['color']
            card.line.width = Pt(3)

            # Faz başlığı
            phase_txt = s.shapes.add_textbox(Inches(1.8), y + Inches(0.12), Inches(10), Inches(0.35))  # CHANGED: daha kompakt başlık
            phase_txt.text_frame.paragraphs[0].text = step['phase']
            phase_txt.text_frame.paragraphs[0].font.size = Pt(22)  # CHANGED: 20->22 daha belirgin
            phase_txt.text_frame.paragraphs[0].font.bold = True
            phase_txt.text_frame.paragraphs[0].font.color.rgb = step['color']

            # İçerik
            item_y_start = y + Inches(0.52)  # CHANGED: başlık altında daha fazla boşluk
            for i, item in enumerate(step['items']):
                item_txt = s.shapes.add_textbox(Inches(2.0), item_y_start + Inches(i * 0.23), Inches(11.8), Inches(0.22))  # CHANGED: daha geniş textbox
                item_txt.text_frame.word_wrap = True  # ADDED: word wrap aktif
                item_txt.text_frame.paragraphs[0].text = f"✓ {item}"
                item_txt.text_frame.paragraphs[0].font.size = Pt(18)  # CHANGED: 20->18 daha kompakt
                item_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            y += Inches(1.6)  # CHANGED: 1.8"->1.6" daha kompakt spacing, total: 2.2+(4×1.6)=8.6" ✓

        # 18. Başarılı Proje Örnekleri (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'success_stories', 0.35)
        self.add_time_badge(s, "11:16")
        self.add_title(s, "⭐ BAŞARILI PROJE ÖRNEKLERİ", 1, 44)

        success_projects = [
            {'title': 'Tarım IoT Sensörü', 'fund': 'TÜBİTAK 2209-A', 'amount': '7.5K TL', 'result': '15 çiftlikte pilot, patent başvurusu', 'team': 'Ziraat Müh. (3 kişi)'},
            {'title': 'Yapay Zeka Radyoloji Asistanı', 'fund': 'KOSGEB Ar-Ge', 'amount': '450K TL', 'result': '2 hastanede kullanımda, %92 doğruluk', 'team': 'Bilgisayar Müh. + Tıp'},
            {'title': 'Akıllı Enerji Yönetimi', 'fund': 'Teknofest 2023', 'amount': '100K TL', 'result': '500 ev/işyerinde %30 tasarruf', 'team': 'Elektrik Müh. (4 kişi)'},
            {'title': 'Eğitim AR Platformu', 'fund': 'Erasmus+ KA2', 'amount': '120K EUR', 'result': '5 ülke, 20 okul, 3000+ öğrenci', 'team': 'Uluslararası ortaklık'},
        ]

        y = Inches(2.5)  # CHANGED from 3.3" to 2.5" - kutuları yukarı al (Eğitim AR Platformu aşağı kalmaması için)
        for proj in success_projects:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(1.5))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['success']
            card.line.width = Pt(3)

            # Proje başlığı
            title_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.15), Inches(7), Inches(0.4))
            title_txt.text_frame.paragraphs[0].text = f"🏆 {proj['title']}"
            title_txt.text_frame.paragraphs[0].font.size = Pt(20)
            title_txt.text_frame.paragraphs[0].font.bold = True
            title_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Fon bilgisi
            fund_txt = s.shapes.add_textbox(Inches(9.5), y + Inches(0.15), Inches(4.5), Inches(0.4))
            fund_txt.text_frame.paragraphs[0].text = f"{proj['fund']} • {proj['amount']}"
            fund_txt.text_frame.paragraphs[0].font.size = Pt(20)
            fund_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

            # Ekip
            team_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.6), Inches(12), Inches(0.3))
            team_txt.text_frame.paragraphs[0].text = f"👥 Ekip: {proj['team']}"
            team_txt.text_frame.paragraphs[0].font.size = Pt(20)
            team_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            # Sonuç
            result_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.95), Inches(12), Inches(0.4))
            result_txt.text_frame.word_wrap = True
            result_txt.text_frame.paragraphs[0].text = f"📊 Sonuç: {proj['result']}"
            result_txt.text_frame.paragraphs[0].font.size = Pt(20)
            result_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['success']

            y += Inches(1.65)

        # 19. Fonlama Takvimi (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'calendar', 0.35)
        self.add_time_badge(s, "11:18")
        self.add_title(s, "📅 2025 FONLAMA TAKVİMİ", 1, 44)

        calendar_items = [
            {'month': 'OCAK-ŞUBAT', 'events': ['TÜBİTAK 2209-A başvuruları', 'KOSGEB 1. dönem başvuruları'], 'color': self.colors['primary']},
            {'month': 'MART-NİSAN', 'events': ['TÜBİTAK değerlendirme', 'Erasmus+ KA2 başvuruları', 'Teknofest ön başvuru'], 'color': self.colors['secondary']},
            {'month': 'MAYIS-HAZİRAN', 'events': ['TÜBİTAK sonuçları', 'KOSGEB değerlendirme', 'Teknofest final başvuru'], 'color': self.colors['success']},
            {'month': 'TEMMUZ-AĞUSTOS', 'events': ['Teknofest yarışmaları', 'Proje başlangıçları', 'Yaz dönemi çalışmaları'], 'color': self.colors['warning']},
            {'month': 'EYLÜL-EKİM', 'events': ['KOSGEB 2. dönem başvuruları', 'Horizon Europe çağrıları', 'TÜBİTAK 1512 başvuruları'], 'color': self.colors['purple']},
            {'month': 'KASIM-ARALIK', 'events': ['Ara raporlar', 'Bütçe revizyonları', 'Yıl sonu sunumları'], 'color': self.colors['pink']},
        ]

        cols = 2
        rows = 3
        x_start = 1.5
        y_start = 2.5  # CHANGED from 2.8" to 2.5" - yukarı alındı
        w = 6.2
        h = 2.0  # CHANGED from 1.8" to 2.0" - event'ler için daha fazla yer
        x_gap = 0.6
        y_gap = 0.2  # CHANGED from 0.3" to 0.2" - daha kompakt

        for i, item in enumerate(calendar_items):
            row = i // cols
            col = i % cols
            x = x_start + col * (w + x_gap)
            y = y_start + row * (h + y_gap)

            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(25, 25, 25)
            card.fill.transparency = 0.2
            card.line.color.rgb = item['color']
            card.line.width = Pt(3)

            # Ay başlığı
            month_txt = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.2), Inches(w - 0.6), Inches(0.4))
            month_txt.text_frame.paragraphs[0].text = item['month']
            month_txt.text_frame.paragraphs[0].font.size = Pt(20)
            month_txt.text_frame.paragraphs[0].font.bold = True
            month_txt.text_frame.paragraphs[0].font.color.rgb = item['color']
            month_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Etkinlikler - TEXTBOX HEIGHT INCREASED (h - 0.9 to h - 1.0)
            events_txt = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.7), Inches(w - 0.6), Inches(h - 1.0))
            events_txt.text_frame.word_wrap = True
            for j, event in enumerate(item['events']):
                if j > 0:
                    p = events_txt.text_frame.add_paragraph()
                else:
                    p = events_txt.text_frame.paragraphs[0]
                p.text = f"• {event}"
                p.font.size = Pt(20)
                p.font.color.rgb = self.colors['white']
                p.space_before = Pt(12) if j > 0 else Pt(0)  # CHANGED from Pt(20) to Pt(12) - daha az boşluk

        # ===== MOLA (11:20-11:30) =====
        print("📍 BÖLÜM 5: Atölye Hazırlık...")

        # 15. Mola
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'brainstorm', 0.5)
        self.add_time_badge(s, "11:20-11:30", x=6)
        self.add_title(s, "☕ MOLA & HAZIRLIK", 2.5, 64, 'pink')

        tasks = [
            '✅ Kısa mola yapın',
            '✅ Laptop/tablet hazırlayın',
            '✅ Claude hesabınızı açın',
            '✅ Yan taraftaki kişiyle tanışın',
            '✅ Takım numaranızı öğrenin',
        ]

        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(4.5), Inches(10), Inches(3.5))  # CHANGED height 3.0"->3.5" (Takım numarası için)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 30, 30)
        box.fill.transparency = 0.2
        box.line.color.rgb = self.colors['pink']
        box.line.width = Pt(4)

        txt = s.shapes.add_textbox(Inches(3.5), Inches(4.8), Inches(9), Inches(2.9))  # CHANGED height 2.4"->2.9"
        txt.text_frame.word_wrap = True
        for i, task in enumerate(tasks):
            if i > 0:
                p = txt.text_frame.add_paragraph()
            else:
                p = txt.text_frame.paragraphs[0]
            p.text = task
            p.font.size = Pt(22)
            p.font.color.rgb = self.colors['white']
            p.space_before = Pt(20)

        # 16. Hesap Kurulum Adımları (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'account_setup', 0.35)
        self.add_time_badge(s, "11:25")
        self.add_title(s, "🔐 HESAP KURULUM REHBERİ", 1, 44)

        setup_steps = [
            {
                'service': 'Claude.ai',
                'steps': ['1. claude.ai adresine git', '2. "Sign Up" ile kayıt ol', '3. Email doğrula', '4. Free plan ile başla'],
                'icon': '🤖',
                'color': self.colors['primary']
            },
            {
                'service': 'GitHub',
                'steps': ['1. github.com/signup', '2. Kullanıcı adı seç', '3. Email doğrula', '4. Ücretsiz plan'],
                'icon': '💻',
                'color': self.colors['secondary']
            },
            {
                'service': 'ChatGPT (Opsiyonel)',
                'steps': ['1. chat.openai.com', '2. Google ile giriş', '3. Free tier yeterli', '4. GPT-3.5 kullan'],
                'icon': '🟢',
                'color': self.colors['success']
            },
            {
                'service': 'Google Gemini (Opsiyonel)',
                'steps': ['1. gemini.google.com', '2. Google hesabı ile giriş', '3. Ücretsiz kullan', '4. Gemini Pro erişimi'],
                'icon': '🔷',
                'color': self.colors['warning']
            },
        ]

        cols = 2
        rows = 2
        x_start = 1.5
        y_start = 2.5  # CHANGED from 2.8" to 2.5" - yukarı al
        w = 6.2
        h = 2.6  # CHANGED from 2.2" to 2.6" - 4 adım için daha fazla yer
        x_gap = 0.6
        y_gap = 0.3

        for i, item in enumerate(setup_steps):
            row = i // cols
            col = i % cols
            x = x_start + col * (w + x_gap)
            y = y_start + row * (h + y_gap)

            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = item['color']
            card.line.width = Pt(3)

            # Icon + Service name
            header = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.2), Inches(w - 0.6), Inches(0.5))
            header.text_frame.paragraphs[0].text = f"{item['icon']} {item['service']}"
            header.text_frame.paragraphs[0].font.size = Pt(20)
            header.text_frame.paragraphs[0].font.bold = True
            header.text_frame.paragraphs[0].font.color.rgb = item['color']
            header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Steps
            steps_txt = s.shapes.add_textbox(Inches(x + 0.4), Inches(y + 0.8), Inches(w - 0.8), Inches(h - 1.1))
            steps_txt.text_frame.word_wrap = True
            for j, step in enumerate(item['steps']):
                if j > 0:
                    p = steps_txt.text_frame.add_paragraph()
                else:
                    p = steps_txt.text_frame.paragraphs[0]
                p.text = step
                p.font.size = Pt(20)
                p.font.color.rgb = self.colors['white']
                p.space_before = Pt(6) if j > 0 else Pt(0)

        # ===== ATÖLYE (11:30-12:00) - ÇOK DETAYLI =====
        print("📍 BÖLÜM 6: Atölye Adım Adım (60 dk)...")

        # 16. Atölye başlangıç
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'meeting', 0.3)
        self.add_time_badge(s, "11:30-12:00")
        self.add_title(s, "⚡ HANDS-ON ATÖLYE", 2.5, 68, 'secondary')
        sub = s.shapes.add_textbox(Inches(3), Inches(5), Inches(10), Inches(0.8))
        sub.text_frame.paragraphs[0].text = "30 Dakikada 25 Proje"
        sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        sub.text_frame.paragraphs[0].font.size = Pt(36)
        sub.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

        # 17. Atölye Yol Haritası (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'roadmap', 0.35)
        self.add_time_badge(s, "11:31")
        self.add_title(s, "🗺️ ATÖLYE YOL HARİTASI", 1, 44)

        roadmap = [
            {'time': '11:30-11:40', 'phase': 'FİKİR', 'tasks': ['Takım oluşturma', 'Problem tanımlama', 'Çözüm beyin fırtınası'], 'color': self.colors['primary']},
            {'time': '11:40-11:50', 'phase': 'PROPOSAL', 'tasks': ['AI ile araştırma', 'Proje önerisi yazma', 'Format kontrolü'], 'color': self.colors['secondary']},
            {'time': '11:50-12:00', 'phase': 'PITCH', 'tasks': ['Pitch hazırlama', 'Prova', 'Sunum'], 'color': self.colors['success']},
        ]

        y = Inches(2.5)  # CHANGED from 3.0" to 2.5" - yukarı al (ATÖLYE YOL HARİTASI)
        for i, item in enumerate(roadmap):
            # Timeline arrow
            if i > 0:
                arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(7.5), y - Inches(0.6), Inches(1), Inches(0.4))  # Küçültüldü
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.colors['light']
                arrow.line.fill.background()

            # Phase card - KÜÇÜLTÜLDÜ 1.8"->1.6"
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), y, Inches(12), Inches(1.6))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = item['color']
            card.line.width = Pt(4)

            # Time badge
            time_box = s.shapes.add_textbox(Inches(2.5), y + Inches(0.2), Inches(2.5), Inches(0.4))
            time_box.text_frame.paragraphs[0].text = f"⏰ {item['time']}"
            time_box.text_frame.paragraphs[0].font.size = Pt(20)
            time_box.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

            # Phase name
            phase_box = s.shapes.add_textbox(Inches(5.5), y + Inches(0.15), Inches(7), Inches(0.5))
            phase_box.text_frame.paragraphs[0].text = item['phase']
            phase_box.text_frame.paragraphs[0].font.size = Pt(24)
            phase_box.text_frame.paragraphs[0].font.bold = True
            phase_box.text_frame.paragraphs[0].font.color.rgb = item['color']

            # Tasks
            tasks_box = s.shapes.add_textbox(Inches(2.5), y + Inches(0.8), Inches(11), Inches(0.9))
            tasks_box.text_frame.word_wrap = True
            tasks_txt = " • ".join(item['tasks'])
            tasks_box.text_frame.paragraphs[0].text = f"📋 {tasks_txt}"
            tasks_box.text_frame.paragraphs[0].font.size = Pt(20)
            tasks_box.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            y += Inches(1.9)  # CHANGED from 2.2" to 1.9" - daha kompakt (1.6" height + 0.3" gap)

        # 18. Design Thinking - Lean - Agile (YENİ! - TAM EKRAN GÖRSEL)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        # Tam ekran görsel - üzerine hiçbir şey ekleme, direkt görseli göster
        try:
            headers = {'User-Agent': 'Mozilla/5.0'}
            r = requests.get(self.images['design_thinking'], timeout=15, headers=headers)
            r.raise_for_status()
            from PIL import Image
            from io import BytesIO
            img = Image.open(BytesIO(r.content))

            # Görseli tam ekrana sığdır
            buf = BytesIO()
            img.save(buf, format='PNG')
            buf.seek(0)

            # Tam ekran görsel ekle - karartma yok, overlay yok, başlık yok!
            left = Inches(0)
            top = Inches(0)
            width = self.prs.slide_width
            height = self.prs.slide_height

            pic = s.shapes.add_picture(buf, left, top, width, height)

            # Progress bar ekle (sadece bunu ekle)
            if hasattr(self, 'slide_counter') and hasattr(self, 'total_slides'):
                self.add_progress_bar(s, self.slide_counter, self.total_slides)

            print(f"  ✅ Design Thinking görseli tam ekran eklendi\n")
        except Exception as e:
            print(f"  ❌ Design Thinking görseli eklenemedi: {str(e)}\n")
            # Fallback - normal arka plan görseli
            self.add_bg_image(s, 'brainstorm', 0.3)
            self.add_title(s, "🎨 DESIGN THINKING - LEAN - AGILE", 0.8, 44, 'warning')

        # 19. Takım Oluşturma Kriterleri (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'team_building', 0.35)
        self.add_time_badge(s, "11:32")
        self.add_title(s, "👥 TAKIM OLUŞTURMA REHBERİ", 1, 44)

        team_criteria = [
            {'icon': '🎯', 'title': 'Tamamlayıcı Beceriler', 'desc': 'Yazılım + Donanım + Tasarım + İş Geliştirme kombinasyonu ideal', 'do': '✅ Farklı bölümlerden', 'dont': '❌ Hepsi aynı alandan'},
            {'icon': '⚖️', 'title': 'Dengeli İş Yükü', 'desc': '2-3 kişi ideal. 4+ olursa koordinasyon zorlaşır', 'do': '✅ 2-3 kişi', 'dont': '❌ Tek kişi veya 5+ kişi'},
            {'icon': '🕐', 'title': 'Zaman Uyumu', 'desc': 'Herkes eşit zaman ayırabilmeli, proje süresince ulaşılabilir olmalı', 'do': '✅ Düzenli toplantı', 'dont': '❌ Farklı şehirler'},
            {'icon': '💡', 'title': 'Ortak Hedef', 'desc': 'Sadece not için değil, gerçekten çözmek istedikleri bir problem olmalı', 'do': '✅ Tutkulu', 'dont': '❌ Sadece geçmek için'},
        ]

        y = Inches(2.5)  # CHANGED to 2.5" - az aşağı al (title çok yakın)
        for item in team_criteria:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(1.1))  # CHANGED height from 1.25" to 1.1" - daha basık
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['secondary']
            card.line.width = Pt(3)

            # Icon - adjusted for smaller card height
            icon_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.25), Inches(0.6), Inches(0.6))
            icon_txt.text_frame.paragraphs[0].text = item['icon']
            icon_txt.text_frame.paragraphs[0].font.size = Pt(28)

            # Title
            title_txt = s.shapes.add_textbox(Inches(2.8), y + Inches(0.1), Inches(4), Inches(0.35))
            title_txt.text_frame.paragraphs[0].text = item['title']
            title_txt.text_frame.paragraphs[0].font.size = Pt(20)
            title_txt.text_frame.paragraphs[0].font.bold = True
            title_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Description
            desc_txt = s.shapes.add_textbox(Inches(2.8), y + Inches(0.48), Inches(7), Inches(0.55))
            desc_txt.text_frame.word_wrap = True
            desc_txt.text_frame.paragraphs[0].text = item['desc']
            desc_txt.text_frame.paragraphs[0].font.size = Pt(20)
            desc_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            # Do/Don't
            tips_txt = s.shapes.add_textbox(Inches(10.5), y + Inches(0.25), Inches(3.5), Inches(0.7))
            tips_txt.text_frame.word_wrap = True
            tips_txt.text_frame.paragraphs[0].text = item['do']
            tips_txt.text_frame.paragraphs[0].font.size = Pt(20)
            tips_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['success']
            p = tips_txt.text_frame.add_paragraph()
            p.text = item['dont']
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['danger']
            p.space_before = Pt(4)

            y += Inches(1.2)  # CHANGED from 1.35" to 1.2" - daha kompakt spacing

        # 34. Takım Oluşturma - SLIDE 34
        # y_start changed from 3.3" to 2.5" to be closer to title

        # 19. Problem Seçim Kriterleri (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'problem_criteria', 0.35)
        self.add_time_badge(s, "11:33")
        self.add_title(s, "🔍 İYİ PROBLEM SEÇİMİ", 1, 44)

        criteria = [
            {'criterion': 'GERÇEKÇİ', 'question': 'Gerçekten var mı?', 'good': 'Fabrikalarda %40 enerji israfı', 'bad': 'İnsanlar teleport olamaması'},
            {'criterion': 'ÖLÇÜLEBİLİR', 'question': 'Sayısal veri var mı?', 'good': 'Yılda 2M TL kayıp', 'bad': '"Bazı insanlar mutsuz"'},
            {'criterion': 'ÇÖZÜMÜ ÇETİN', 'question': 'Şu an çözülmemiş mi?', 'good': 'Manuel takip çok zaman alıyor', 'bad': 'Excel yeterli'},
            {'criterion': 'KAPSAM', 'question': '6-12 ayda bitirilebilir mi?', 'good': 'Prototip + 10 pilot test', 'bad': 'Dünyayı değiştirme'},
            {'criterion': 'UZMANLAŞTIRILABİLİR', 'question': 'Ekibin bilgisi yeterli mi?', 'good': 'Bilgisayar müh. + Endüstri müh.', 'bad': 'Nükleer füzyon'},
        ]

        y = Inches(2.5)  # CHANGED to 2.5" - az aşağı al (title çok yakın)
        for item in criteria:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y, Inches(14), Inches(0.95))  # CHANGED height from 1.0" to 0.95" - biraz daha basık
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(25, 25, 25)
            card.fill.transparency = 0.2
            card.line.color.rgb = self.colors['primary']
            card.line.width = Pt(2)

            # Criterion label
            crit_txt = s.shapes.add_textbox(Inches(1.3), y + Inches(0.25), Inches(2.5), Inches(0.5))
            crit_txt.text_frame.paragraphs[0].text = item['criterion']
            crit_txt.text_frame.paragraphs[0].font.size = Pt(20)
            crit_txt.text_frame.paragraphs[0].font.bold = True
            crit_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

            # Question
            q_txt = s.shapes.add_textbox(Inches(4), y + Inches(0.25), Inches(3), Inches(0.5))
            q_txt.text_frame.paragraphs[0].text = item['question']
            q_txt.text_frame.paragraphs[0].font.size = Pt(20)
            q_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Good example
            good_txt = s.shapes.add_textbox(Inches(7.5), y + Inches(0.15), Inches(7.2), Inches(0.35))  # CHANGED width 3.5"->7.2" (daha geniş, 15"a kadar)
            good_txt.text_frame.word_wrap = True
            good_txt.text_frame.paragraphs[0].text = f"✅ {item['good']}"
            good_txt.text_frame.paragraphs[0].font.size = Pt(20)
            good_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['success']

            # Bad example
            bad_txt = s.shapes.add_textbox(Inches(7.5), y + Inches(0.55), Inches(7.2), Inches(0.35))  # CHANGED width 3.5"->7.2" (daha geniş, 15"a kadar)
            bad_txt.text_frame.word_wrap = True
            bad_txt.text_frame.paragraphs[0].text = f"❌ {item['bad']}"
            bad_txt.text_frame.paragraphs[0].font.size = Pt(20)
            bad_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['danger']

            y += Inches(1.1)  # CHANGED from 1.05" to 1.1" - daha fazla boşluk (0.15" gap)

        # 20. Problem Örnekleri (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'problem_examples', 0.35)
        self.add_time_badge(s, "11:34")
        self.add_title(s, "💡 PROBLEM ÖRNEKLERİ", 1, 44)

        problems = [
            {'domain': 'TARIM', 'problem': 'Küçük çiftçiler toprak nemini manuel ölçüyor', 'solution': 'IoT sensörlü akıllı sulama', 'impact': '%40 su tasarrufu'},
            {'domain': 'SAĞLIK', 'problem': 'Kronik hasta ilaç saatlerini unutuyor', 'solution': 'Akıllı hatırlatma + takip uygulaması', 'impact': '%60 uyum artışı'},
            {'domain': 'EĞİTİM', 'problem': 'Uzaktan eğitimde öğrenci motivasyonu düşük', 'solution': 'Gamification + AI mentor', 'impact': '%35 tamamlama oranı'},
            {'domain': 'LOJİSTİK', 'problem': 'Şehir içi teslimat rota optimizasyonu zayıf', 'solution': 'ML tabanlı dinamik rotalama', 'impact': '%25 yakıt tasarrufu'},
            {'domain': 'ENERJİ', 'problem': 'Binalarda gereksiz aydınlatma/ısıtma', 'solution': 'Doluluk sensörlü akıllı otomasyon', 'impact': '%30 enerji düşüşü'},
            {'domain': 'ATIK', 'problem': 'Geri dönüşüm kutularına yanlış atık', 'solution': 'Görüntü tanıma ile akıllı sınıflandırma', 'impact': '%50 geri dönüşüm'},
        ]

        cols = 2
        rows = 3
        x_start = 1.5
        y_start = 2.3  # CHANGED to 2.3" - az aşağı al için yukarı başla
        w = 6.2
        h = 1.75  # CHANGED from 1.65" to 1.75" - aşağı doğru büyüt
        x_gap = 0.6
        y_gap = 0.3  # CHANGED from 0.25" to 0.3" - daha fazla boşluk

        for i, item in enumerate(problems):
            row = i // cols
            col = i % cols
            x = x_start + col * (w + x_gap)
            y = y_start + row * (h + y_gap)

            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['teal']
            card.line.width = Pt(3)

            # Domain badge
            domain_txt = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.15), Inches(w - 0.6), Inches(0.35))
            domain_txt.text_frame.paragraphs[0].text = f"🏷️ {item['domain']}"
            domain_txt.text_frame.paragraphs[0].font.size = Pt(20)
            domain_txt.text_frame.paragraphs[0].font.bold = True
            domain_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

            # Problem
            prob_txt = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.55), Inches(w - 0.6), Inches(0.4))
            prob_txt.text_frame.word_wrap = True
            prob_txt.text_frame.paragraphs[0].text = f"❗ {item['problem']}"
            prob_txt.text_frame.paragraphs[0].font.size = Pt(20)
            prob_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Solution
            sol_txt = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.95), Inches(w - 0.6), Inches(0.3))
            sol_txt.text_frame.word_wrap = True
            sol_txt.text_frame.paragraphs[0].text = f"💡 {item['solution']}"
            sol_txt.text_frame.paragraphs[0].font.size = Pt(20)
            sol_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

            # Impact
            impact_txt = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.3), Inches(w - 0.6), Inches(0.25))
            impact_txt.text_frame.paragraphs[0].text = f"📈 {item['impact']}"
            impact_txt.text_frame.paragraphs[0].font.size = Pt(20)
            impact_txt.text_frame.paragraphs[0].font.bold = True
            impact_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['success']

        # 21. SWOT Analizi Şablonu (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'swot_analysis', 0.35)
        self.add_time_badge(s, "11:35")
        self.add_title(s, "📊 SWOT ANALİZİ ŞABLONU", 1, 44)

        swot_items = [
            {'type': 'STRENGTHS (Güçlü Yönler)', 'questions': ['Takımın uzmanlığı nedir?', 'Hangi kaynaklara erişim var?', 'Benzersiz avantajımız nedir?'], 'color': self.colors['success']},
            {'type': 'WEAKNESSES (Zayıf Yönler)', 'questions': ['Hangi becerilerde eksiklik var?', 'Bütçe/zaman kısıtları nedir?', 'Hangi riskleri alamayız?'], 'color': self.colors['danger']},
            {'type': 'OPPORTUNITIES (Fırsatlar)', 'questions': ['Hangi trendler işimize yarar?', 'Ortaklık fırsatları?', 'Pazar boşlukları?'], 'color': self.colors['primary']},
            {'type': 'THREATS (Tehditler)', 'questions': ['Rakipler kimler?', 'Teknolojik riskler?', 'Regülasyon engelleri?'], 'color': self.colors['warning']},
        ]

        positions = [
            {'x': 1.5, 'y': 2.8},  # Top-left
            {'x': 8.3, 'y': 2.8},  # Top-right
            {'x': 1.5, 'y': 5.5},  # Bottom-left
            {'x': 8.3, 'y': 5.5},  # Bottom-right
        ]

        for i, item in enumerate(swot_items):
            pos = positions[i]
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(pos['x']), Inches(pos['y']), Inches(6.2), Inches(2.4))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = item['color']
            card.line.width = Pt(4)

            # Type header
            type_txt = s.shapes.add_textbox(Inches(pos['x'] + 0.3), Inches(pos['y'] + 0.2), Inches(5.6), Inches(0.4))
            type_txt.text_frame.paragraphs[0].text = item['type']
            type_txt.text_frame.paragraphs[0].font.size = Pt(20)
            type_txt.text_frame.paragraphs[0].font.bold = True
            type_txt.text_frame.paragraphs[0].font.color.rgb = item['color']
            type_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Questions
            questions_txt = s.shapes.add_textbox(Inches(pos['x'] + 0.4), Inches(pos['y'] + 0.7), Inches(5.4), Inches(1.5))
            questions_txt.text_frame.word_wrap = True
            for j, q in enumerate(item['questions']):
                if j > 0:
                    p = questions_txt.text_frame.add_paragraph()
                else:
                    p = questions_txt.text_frame.paragraphs[0]
                p.text = f"• {q}"
                p.font.size = Pt(20)
                p.font.color.rgb = self.colors['white']
                p.space_before = Pt(6) if j > 0 else Pt(0)

        # 22. Çözüm Matrisi (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'solution_matrix', 0.35)
        self.add_time_badge(s, "11:36")
        self.add_title(s, "🎯 ÇÖZÜM DEĞERLENDİRME MATRİSİ", 1, 40)

        # Matrix header
        header_y = Inches(3.0)
        headers = ['Çözüm Fikri', 'Teknik\nZorluk (1-5)', 'Maliyet\n(1-5)', 'Etki\n(1-5)', 'Toplam\nSkor']
        header_widths = [4.5, 2, 2, 2, 2]
        header_x = Inches(1.5)

        for i, header in enumerate(headers):
            h_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, header_x, header_y, Inches(header_widths[i]), Inches(0.6))
            h_box.fill.solid()
            h_box.fill.fore_color.rgb = self.colors['primary']
            h_box.line.fill.background()

            h_txt = s.shapes.add_textbox(header_x, header_y, Inches(header_widths[i]), Inches(0.6))
            h_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            h_txt.text_frame.paragraphs[0].text = header
            h_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            h_txt.text_frame.paragraphs[0].font.size = Pt(20)
            h_txt.text_frame.paragraphs[0].font.bold = True
            h_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            header_x += Inches(header_widths[i])

        # Example rows
        examples = [
            {'solution': 'IoT Sensörlü Akıllı Sulama', 'tech': 3, 'cost': 2, 'impact': 5, 'total': 10},
            {'solution': 'Blockchain Tabanlı Tedarik Zinciri', 'tech': 5, 'cost': 4, 'impact': 4, 'total': 5},
            {'solution': 'Mobil Uygulama ile İlaç Takibi', 'tech': 2, 'cost': 1, 'impact': 4, 'total': 11},
        ]

        row_y = header_y + Inches(0.7)
        for ex in examples:
            row_x = Inches(1.5)

            # Solution name
            sol_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, row_x, row_y, Inches(header_widths[0]), Inches(0.5))
            sol_box.fill.solid()
            sol_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
            sol_box.fill.transparency = 0.3
            sol_box.line.color.rgb = self.colors['light']
            sol_box.line.width = Pt(1)

            sol_txt = s.shapes.add_textbox(row_x + Inches(0.2), row_y, Inches(header_widths[0] - 0.4), Inches(0.5))
            sol_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            sol_txt.text_frame.paragraphs[0].text = ex['solution']
            sol_txt.text_frame.paragraphs[0].font.size = Pt(20)
            sol_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            row_x += Inches(header_widths[0])

            # Scores
            scores = [ex['tech'], ex['cost'], ex['impact'], ex['total']]
            for i, score in enumerate(scores):
                score_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, row_x, row_y, Inches(header_widths[i+1]), Inches(0.5))
                score_box.fill.solid()
                score_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
                score_box.fill.transparency = 0.3
                score_box.line.color.rgb = self.colors['light']
                score_box.line.width = Pt(1)

                score_txt = s.shapes.add_textbox(row_x, row_y, Inches(header_widths[i+1]), Inches(0.5))
                score_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                score_txt.text_frame.paragraphs[0].text = str(score)
                score_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                score_txt.text_frame.paragraphs[0].font.size = Pt(20)
                score_txt.text_frame.paragraphs[0].font.bold = True
                score_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['success'] if i == 3 and score >= 10 else self.colors['white']

                row_x += Inches(header_widths[i+1])

            row_y += Inches(0.6)

        # Formula explanation
        formula_txt = s.shapes.add_textbox(Inches(1.5), row_y + Inches(0.3), Inches(13), Inches(0.8))
        formula_txt.text_frame.word_wrap = True
        formula_txt.text_frame.paragraphs[0].text = "📐 FORMÜL: Toplam Skor = (6 - Teknik Zorluk) + (6 - Maliyet) + Etki"
        formula_txt.text_frame.paragraphs[0].font.size = Pt(20)
        formula_txt.text_frame.paragraphs[0].font.bold = True
        formula_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']
        p = formula_txt.text_frame.add_paragraph()
        p.text = "🎯 Hedef: En yüksek skoru seçin (10+ mükemmel, 7-9 iyi, <7 riskli)"
        p.font.size = Pt(20)
        p.font.color.rgb = self.colors['light']
        p.space_before = Pt(20)

        # 23. Mentor İpuçları (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'mentor_tips', 0.35)
        self.add_time_badge(s, "11:37")
        self.add_title(s, "🎓 MENTOR İPUÇLARI", 1, 44)

        tips = [
            {'phase': 'FİKİR AŞAMASI', 'tips': ['Geniş düşünün, sonra daraltın', 'Kendinize "Neden?" sorusunu 5 kez sorun', 'Rakipleri araştırın (Google Scholar, Patent DB)'], 'icon': '💡'},
            {'phase': 'ARAŞTIRMA', 'tips': ['DeepResearch ile son 2 yıl makalelerine bakın', 'LinkedIn\'de sektör uzmanlarını bulun', 'Benzer projelerin bütçesini inceleyin'], 'icon': '🔬'},
            {'phase': 'YAZIM', 'tips': ['Her bölümü AI\'ya yazdırın, sonra düzenleyin', 'Sayısal veri olmadan cümle yazmayın', 'Basit dil kullanın, jargon minimumda'], 'icon': '✍️'},
            {'phase': 'SUNUM', 'tips': ['İlk 10 saniye kritik, dikkat çekin', 'Demo video hazırlayın (Loom, OBS)', 'Soruları önceden tahmin edin, cevap hazırlayın'], 'icon': '🎤'},
        ]

        y = Inches(2.5)  # CHANGED from 3.3" to 2.5" - yukarı al title yaklaştır (MENTOR İPUÇLARI)
        for item in tips:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(1.3))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['teal']
            card.line.width = Pt(3)

            # Icon + Phase
            header_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.15), Inches(12), Inches(0.4))
            header_txt.text_frame.paragraphs[0].text = f"{item['icon']} {item['phase']}"
            header_txt.text_frame.paragraphs[0].font.size = Pt(20)
            header_txt.text_frame.paragraphs[0].font.bold = True
            header_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

            # Tips
            tips_txt = s.shapes.add_textbox(Inches(2.3), y + Inches(0.6), Inches(12), Inches(0.6))
            tips_txt.text_frame.word_wrap = True
            tips_txt.text_frame.paragraphs[0].text = " • ".join(item['tips'])
            tips_txt.text_frame.paragraphs[0].font.size = Pt(20)
            tips_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            y += Inches(1.4)

        # 24. Pitch Şablonu (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'pitch_template', 0.35)
        self.add_time_badge(s, "11:38")
        self.add_title(s, "🎬 1-DAKİKALIK PITCH ŞABLONU", 1, 40)

        pitch_structure = [
            {'section': 'HOOK (5 sn)', 'content': 'Dikkat çekici istatistik veya soru', 'example': '"Fabrikaların %40\'ı enerjiyi boşa harcıyor. Peki ya senin okulun?"'},
            {'section': 'PROBLEM (15 sn)', 'content': 'Ne sorunu çözüyorsunuz? Kimleri etkiliyor?', 'example': '"Türkiye\'de yıllık 10B TL enerji israfı var. KOBİler manuel izleme yapıyor, verimsiz."'},
            {'section': 'ÇÖZÜM (20 sn)', 'content': 'Nasıl çözüyorsunuz? Teknoloji nedir?', 'example': '"IoT sensörler + yapay zeka ile gerçek zamanlı enerji optimizasyonu. Otomatik uyarı ve raporlama."'},
            {'section': 'ETKİ (10 sn)', 'content': 'Sonuç ne olur? Sayısal hedef', 'example': '"%30 enerji tasarrufu, 6 ayda ROI, 100 KOBİ\'ye ölçeklenebilir."'},
            {'section': 'CALL-TO-ACTION (10 sn)', 'content': 'Ne istiyorsunuz? Sonraki adım?', 'example': '"TÜBİTAK 2209-A başvurusu yapacağız. Pilot için 3 fabrika arıyoruz."'},
        ]

        y = Inches(3.0)
        for item in pitch_structure:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y, Inches(14), Inches(1))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(25, 25, 25)
            card.fill.transparency = 0.2
            card.line.color.rgb = self.colors['secondary']
            card.line.width = Pt(2)

            # Section name
            sec_txt = s.shapes.add_textbox(Inches(1.3), y + Inches(0.15), Inches(2.5), Inches(0.4))
            sec_txt.text_frame.paragraphs[0].text = item['section']
            sec_txt.text_frame.paragraphs[0].font.size = Pt(20)
            sec_txt.text_frame.paragraphs[0].font.bold = True
            sec_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

            # Content description
            cont_txt = s.shapes.add_textbox(Inches(4), y + Inches(0.15), Inches(4.5), Inches(0.4))
            cont_txt.text_frame.word_wrap = True
            cont_txt.text_frame.paragraphs[0].text = item['content']
            cont_txt.text_frame.paragraphs[0].font.size = Pt(20)
            cont_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            # Example
            ex_txt = s.shapes.add_textbox(Inches(1.3), y + Inches(0.58), Inches(13.2), Inches(0.4))
            ex_txt.text_frame.word_wrap = True
            ex_txt.text_frame.paragraphs[0].text = f"📝 {item['example']}"
            ex_txt.text_frame.paragraphs[0].font.size = Pt(20)
            ex_txt.text_frame.paragraphs[0].font.italic = True
            ex_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['teal']

            y += Inches(1.05)

        # 25. Feedback Loop (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'feedback_loop', 0.35)
        self.add_time_badge(s, "11:39")
        self.add_title(s, "🔄 FEEDBACK VE İYİLEŞTİRME", 1, 44)

        feedback_process = [
            {'step': '1', 'action': 'Takım İçi Review', 'method': 'Her bölümü sesli okuyun, anlaşılmayan yer var mı?', 'time': '5 dk'},
            {'step': '2', 'action': 'Peer Review (Yan Takım)', 'method': '3-2-1 Feedback: 3 güçlü yön, 2 iyileştirme, 1 soru', 'time': '10 dk'},
            {'step': '3', 'action': 'AI Review', 'method': 'Claude\'a tüm metni verin: "Bu proje önerisini değerlendir, eksikleri söyle"', 'time': '3 dk'},
            {'step': '4', 'action': 'Revizyon', 'method': 'Geri bildirimleri entegre edin, tekrar AI ile kontrol', 'time': '10 dk'},
        ]

        y = Inches(3)
        for item in feedback_process:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), y, Inches(12), Inches(1.3))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['primary']
            card.line.width = Pt(3)

            # Step number circle
            circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), y + Inches(0.35), Inches(0.6), Inches(0.6))
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors['primary']
            circle.line.fill.background()

            num_txt = s.shapes.add_textbox(Inches(2.4), y + Inches(0.35), Inches(0.6), Inches(0.6))
            num_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            num_txt.text_frame.paragraphs[0].text = item['step']
            num_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_txt.text_frame.paragraphs[0].font.size = Pt(20)
            num_txt.text_frame.paragraphs[0].font.bold = True
            num_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Action
            action_txt = s.shapes.add_textbox(Inches(3.2), y + Inches(0.15), Inches(5), Inches(0.4))
            action_txt.text_frame.paragraphs[0].text = item['action']
            action_txt.text_frame.paragraphs[0].font.size = Pt(20)
            action_txt.text_frame.paragraphs[0].font.bold = True
            action_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Time badge
            time_txt = s.shapes.add_textbox(Inches(12), y + Inches(0.45), Inches(1.5), Inches(0.4))
            time_txt.text_frame.paragraphs[0].text = f"⏱️ {item['time']}"
            time_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            time_txt.text_frame.paragraphs[0].font.size = Pt(20)
            time_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

            # Method
            method_txt = s.shapes.add_textbox(Inches(3.2), y + Inches(0.65), Inches(10.3), Inches(0.5))
            method_txt.text_frame.word_wrap = True
            method_txt.text_frame.paragraphs[0].text = f"📋 {item['method']}"
            method_txt.text_frame.paragraphs[0].font.size = Pt(20)
            method_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            y += Inches(1.4)

        # 26. Template Demo (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'template_demo', 0.35)
        self.add_time_badge(s, "11:40")
        self.add_title(s, "📄 PROJE ŞABLONU ERİŞİMİ", 1, 44)

        templates = [
            {'name': 'TÜBİTAK 2209-A Şablonu', 'format': 'Word (.docx)', 'pages': 'Max 20 sayfa', 'link': 'tubitak.gov.tr/tr/burslar/lisans/burs-programlari/2209-a', 'desc': 'Resmi TÜBİTAK başvuru formu, Arial 9 font'},
            {'name': 'KOSGEB Rize Rehberi', 'format': 'Web Rehber', 'pages': 'Detaylı kılavuz', 'link': 'coruhtech.github.io/kosgeb-rize/', 'desc': 'KOSGEB başvuru süreçleri ve iş planı hazırlama rehberi'},
            {'name': 'Pitch Deck Şablonu', 'format': 'PowerPoint', 'pages': '10 slayt', 'link': 'canva.com/presentations/templates/pitch-deck/', 'desc': 'Ücretsiz Canva şablonları, özelleştirilebilir'},
            {'name': 'GitHub Proje Şablonları', 'format': 'Markdown', 'pages': 'README', 'link': 'github.com/ucoruh', 'desc': 'Açık kaynak proje dokümantasyon örnekleri'},
        ]

        y = Inches(3)
        for item in templates:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(1.2))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['success']
            card.line.width = Pt(3)

            # Template name
            name_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.15), Inches(6), Inches(0.4))
            name_txt.text_frame.paragraphs[0].text = f"📄 {item['name']}"
            name_txt.text_frame.paragraphs[0].font.size = Pt(20)
            name_txt.text_frame.paragraphs[0].font.bold = True
            name_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Format + Pages
            meta_txt = s.shapes.add_textbox(Inches(8.5), y + Inches(0.15), Inches(5), Inches(0.4))
            meta_txt.text_frame.paragraphs[0].text = f"{item['format']} • {item['pages']}"
            meta_txt.text_frame.paragraphs[0].font.size = Pt(20)
            meta_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

            # Description
            desc_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.6), Inches(8), Inches(0.5))
            desc_txt.text_frame.word_wrap = True
            desc_txt.text_frame.paragraphs[0].text = item['desc']
            desc_txt.text_frame.paragraphs[0].font.size = Pt(20)
            desc_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            # Link
            link_txt = s.shapes.add_textbox(Inches(10.5), y + Inches(0.65), Inches(3.5), Inches(0.4))
            link_txt.text_frame.paragraphs[0].text = f"🔗 {item['link']}"
            link_txt.text_frame.paragraphs[0].font.size = Pt(20)
            link_txt.text_frame.paragraphs[0].font.underline = True
            link_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

            y += Inches(1.3)

        # 27. Research Methods (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'research_methods', 0.35)
        self.add_time_badge(s, "11:41")
        self.add_title(s, "🔬 ARAŞTIRMA YÖNTEMLERİ", 1, 44)

        research_tools = [
            {'tool': 'Claude DeepResearch', 'use': 'Literatür tarama', 'how': 'En son makaleleri, trendleri bulur', 'time': '3 dk', 'icon': '🤖'},
            {'tool': 'Google Scholar', 'use': 'Akademik makaleler', 'how': 'Anahtar kelime ara, son 2 yıl filtrele, atıf sayısına bak', 'time': '10 dk', 'icon': '📚'},
            {'tool': 'Patent Databases', 'use': 'Benzer patentler', 'how': 'espacenet.com, USPTO.gov, yenilik kontrolü', 'time': '15 dk', 'icon': '⚖️'},
            {'tool': 'Statista / İTO', 'use': 'Pazar verileri', 'how': 'Sektör büyüklüğü, trend analizi, Türkiye istatistikleri', 'time': '10 dk', 'icon': '📊'},
            {'tool': 'LinkedIn Uzman Görüşü', 'use': 'Sektör doğrulaması', 'how': 'İlgili uzmanlara mesaj, 3-4 soruya cevap iste', 'time': '1-2 gün', 'icon': '💼'},
        ]

        y = Inches(2.5)  # CHANGED from 3.3" to 2.5" - yukarı al (ARAŞTIRMA YÖNTEMLERİ)
        for item in research_tools:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(13), Inches(1))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['teal']
            card.line.width = Pt(3)

            # Icon + Tool name
            header_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.1), Inches(4), Inches(0.35))
            header_txt.text_frame.paragraphs[0].text = f"{item['icon']} {item['tool']}"
            header_txt.text_frame.paragraphs[0].font.size = Pt(20)
            header_txt.text_frame.paragraphs[0].font.bold = True
            header_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Use case
            use_txt = s.shapes.add_textbox(Inches(6.5), y + Inches(0.1), Inches(2.5), Inches(0.35))
            use_txt.text_frame.paragraphs[0].text = f"🎯 {item['use']}"
            use_txt.text_frame.paragraphs[0].font.size = Pt(20)
            use_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

            # Time
            time_txt = s.shapes.add_textbox(Inches(12.5), y + Inches(0.2), Inches(1.5), Inches(0.35))
            time_txt.text_frame.paragraphs[0].text = f"⏱️ {item['time']}"
            time_txt.text_frame.paragraphs[0].font.size = Pt(20)
            time_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['light']
            time_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # How to use
            how_txt = s.shapes.add_textbox(Inches(2), y + Inches(0.55), Inches(12), Inches(0.4))
            how_txt.text_frame.word_wrap = True
            how_txt.text_frame.paragraphs[0].text = f"💡 {item['how']}"
            how_txt.text_frame.paragraphs[0].font.size = Pt(20)
            how_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            y += Inches(1.1)

        # 28. FAZ 1: Fikir Geliştirme (11:30-11:40)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'innovation', 0.3)
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(1.2), Inches(1.2))
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.colors['secondary']
        badge.line.fill.background()
        num = s.shapes.add_textbox(Inches(1), Inches(1), Inches(1.2), Inches(1.2))
        num.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        num.text_frame.paragraphs[0].text = "FAZ\n1"
        num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        num.text_frame.paragraphs[0].font.size = Pt(20)
        num.text_frame.paragraphs[0].font.bold = True
        num.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

        title = s.shapes.add_textbox(Inches(2.5), Inches(1), Inches(12), Inches(0.8))
        title.text_frame.paragraphs[0].text = "FİKİR GELİŞTİRME"
        title.text_frame.paragraphs[0].font.size = Pt(42)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

        time = s.shapes.add_textbox(Inches(2.5), Inches(1.9), Inches(8), Inches(0.5))
        time.text_frame.paragraphs[0].text = "⏰ 11:30-11:40 (10 dk)"
        time.text_frame.paragraphs[0].font.size = Pt(20)
        time.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3), Inches(13), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(25, 25, 25)
        box.fill.transparency = 0.2
        box.line.color.rgb = self.colors['secondary']
        box.line.width = Pt(3)

        steps = [
            '👉 2\'şerli takımlar oluşun',
            '👉 3 problem belirleyin (gerçek sorunlar)',
            '👉 Claude\'a sorun: Her problem için 3 çözüm',
            '👉 En uygulanabilir çözümü seçin',
            '👉 Yan takımla paylaşın (3-2-1 feedback)',
            '',
            '✅ ÇIKTI: 1 sayfa problem-çözüm açıklaması'
        ]

        txt = s.shapes.add_textbox(Inches(2), Inches(3.3), Inches(12), Inches(3.9))
        txt.text_frame.word_wrap = True
        for i, step in enumerate(steps):
            if i > 0:
                p = txt.text_frame.add_paragraph()
            else:
                p = txt.text_frame.paragraphs[0]
            p.text = step
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['white']
            p.space_before = Pt(20)
            if step.startswith('✅'):
                p.font.bold = True
                p.font.color.rgb = self.colors['success']

        # 18. PROJE ÖNERİSİ FORMATI - TÜBİTAK ŞABLONU (YENİ!)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'document', 0.3)
        self.add_time_badge(s, "11:35")
        self.add_title(s, "📝 PROJE ÖNERİSİ FORMATI", 0.8, 40, 'warning')

        subtitle = s.shapes.add_textbox(Inches(2), Inches(2.1), Inches(12), Inches(0.4))  # CHANGED: 1.8"->2.1" aşağı al
        subtitle.text_frame.paragraphs[0].text = "TÜBİTAK 2209-A Başvuru Şablonu"
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.italic = True
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

        format_items = [
            {'num': '1', 'section': 'PROJE BAŞLIĞI', 'detail': 'Kısa, çekici, açıklayıcı (max 15 kelime)', 'ex': 'Örn: "Akıllı Sensörlerle Endüstriyel Enerji Verimliliği"'},
            {'num': '2', 'section': 'PROBLEM TANIMI', 'detail': 'Hangi sorunu çözüyor? Nicel verilerle', 'ex': 'Örn: "Fabrikalarda %40 enerji israfı, yıllık 2M TL kayıp"'},
            {'num': '3', 'section': 'ÇÖZÜM ÖNERİSİ', 'detail': 'Nasıl çözülecek? Teknoloji + Yöntem', 'ex': 'Örn: "IoT sensörler + ML ile gerçek zamanlı optimizasyon"'},
            {'num': '4', 'section': 'YENİLİK & ÖZGÜNLÜK', 'detail': 'Mevcut çözümlerden farkı nedir?', 'ex': 'Örn: "Halihazırda manuel takip var, bizim sistem otonom"'},
        ]

        y = Inches(3.0)
        for item in format_items:
            # Ana kart
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y, Inches(14), Inches(1.2))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(30, 30, 30)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['primary']
            card.line.width = Pt(2)

            # Numara rozeti
            num_badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.3), y + Inches(0.3), Inches(0.6), Inches(0.6))
            num_badge.fill.solid()
            num_badge.fill.fore_color.rgb = self.colors['warning']
            num_badge.line.fill.background()

            num_txt = s.shapes.add_textbox(Inches(1.3), y + Inches(0.3), Inches(0.6), Inches(0.6))
            num_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            num_txt.text_frame.paragraphs[0].text = item['num']
            num_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_txt.text_frame.paragraphs[0].font.size = Pt(20)
            num_txt.text_frame.paragraphs[0].font.bold = True
            num_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']

            # Bölüm adı
            sec = s.shapes.add_textbox(Inches(2.2), y + Inches(0.1), Inches(5), Inches(0.4))
            sec.text_frame.paragraphs[0].text = item['section']
            sec.text_frame.paragraphs[0].font.size = Pt(20)
            sec.text_frame.paragraphs[0].font.bold = True
            sec.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Detay
            det = s.shapes.add_textbox(Inches(2.2), y + Inches(0.5), Inches(5.5), Inches(0.5))
            det.text_frame.word_wrap = True
            det.text_frame.paragraphs[0].text = item['detail']
            det.text_frame.paragraphs[0].font.size = Pt(20)
            det.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            # Örnek
            ex = s.shapes.add_textbox(Inches(8), y + Inches(0.3), Inches(6.5), Inches(0.6))
            ex.text_frame.word_wrap = True
            ex.text_frame.paragraphs[0].text = item['ex']
            ex.text_frame.paragraphs[0].font.size = Pt(20)
            ex.text_frame.paragraphs[0].font.italic = True
            ex.text_frame.paragraphs[0].font.color.rgb = self.colors['teal']

            y += Inches(1.3)

        # 19. PROJE ÖNERİSİ FORMATI - DEVAM (YENİ!)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'writing', 0.3)
        self.add_time_badge(s, "11:36")
        self.add_title(s, "📝 PROJE ÖNERİSİ FORMATI (Devam)", 0.8, 36, 'warning')

        format_items2 = [
            {'num': '5', 'section': 'HEDEFLER', 'detail': '3-5 ölçülebilir hedef (SMART)', 'ex': 'Örn: "6 ayda %30 enerji tasarrufu, 10 sensör kurulumu"'},
            {'num': '6', 'section': 'YÖNTEM & İŞ PLANI', 'detail': 'Adım adım nasıl yapılacak? (Gantt)', 'ex': 'Örn: "Ay 1-2: Tasarım, Ay 3-4: Prototip, Ay 5-6: Test"'},
            {'num': '7', 'section': 'BEKLENİLEN ÇIKTILAR', 'detail': 'Sonuçlar somut, ölçülebilir olmalı', 'ex': 'Örn: "Çalışan prototip, test raporu, 1 makale taslağı"'},
            {'num': '8', 'section': 'BÜTÇE', 'detail': 'Kalemler: Malzeme, yazılım, hizmet alımı', 'ex': 'Örn: "Sensörler: 3K, Cloud: 1.5K, Diğer: 3K = 7.5K TL"'},
            {'num': '9', 'section': 'REFERANSLAR', 'detail': 'Bilimsel kaynaklar (min 5 makale)', 'ex': 'Örn: "IEEE IoT Journal, Energy Management 2024..."'},
        ]

        y = Inches(2.3)  # CHANGED from 2.0" to 2.3" - title'dan uzaklaştır (PROJE ÖNERİSİ Devam)
        for item in format_items2:
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y, Inches(14), Inches(1.1))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(30, 30, 30)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['primary']
            card.line.width = Pt(2)

            num_badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.3), y + Inches(0.25), Inches(0.6), Inches(0.6))
            num_badge.fill.solid()
            num_badge.fill.fore_color.rgb = self.colors['warning']
            num_badge.line.fill.background()

            num_txt = s.shapes.add_textbox(Inches(1.3), y + Inches(0.25), Inches(0.6), Inches(0.6))
            num_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            num_txt.text_frame.paragraphs[0].text = item['num']
            num_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_txt.text_frame.paragraphs[0].font.size = Pt(20)
            num_txt.text_frame.paragraphs[0].font.bold = True
            num_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']

            sec = s.shapes.add_textbox(Inches(2.2), y + Inches(0.1), Inches(4.5), Inches(0.4))
            sec.text_frame.paragraphs[0].text = item['section']
            sec.text_frame.paragraphs[0].font.size = Pt(20)
            sec.text_frame.paragraphs[0].font.bold = True
            sec.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            det = s.shapes.add_textbox(Inches(2.2), y + Inches(0.5), Inches(5.5), Inches(0.5))
            det.text_frame.word_wrap = True
            det.text_frame.paragraphs[0].text = item['detail']
            det.text_frame.paragraphs[0].font.size = Pt(20)
            det.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

            ex = s.shapes.add_textbox(Inches(8), y + Inches(0.25), Inches(6.5), Inches(0.6))
            ex.text_frame.word_wrap = True
            ex.text_frame.paragraphs[0].text = item['ex']
            ex.text_frame.paragraphs[0].font.size = Pt(20)
            ex.text_frame.paragraphs[0].font.italic = True
            ex.text_frame.paragraphs[0].font.color.rgb = self.colors['teal']

            y += Inches(1.2)

        # 20. AI İLE PROJE YAZIM ADIMLARI (YENİ!)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'typing', 0.3)
        self.add_time_badge(s, "11:38")
        self.add_title(s, "🤖 AI İLE PROJE NASIL YAZILIR?", 0.8, 40, 'success')

        ai_steps = [
            {
                'step': 'ADIM 1',
                'action': 'Claude\'a problemi tanımlayın',
                'prompt': '"Fabrikalarda enerji israfı problemi hakkında bilgi ver"',
                'time': '2 dk'
            },
            {
                'step': 'ADIM 2',
                'action': 'DeepResearch ile araştırma',
                'prompt': '"IoT enerji izleme sistemleri hakkında güncel makaleler bul"',
                'time': '3 dk'
            },
            {
                'step': 'ADIM 3',
                'action': 'Proje önerisi outline oluştur',
                'prompt': '"TÜBİTAK 2209-A formatında proje önerisi şablonu oluştur"',
                'time': '2 dk'
            },
            {
                'step': 'ADIM 4',
                'action': 'Her bölümü AI ile doldurun',
                'prompt': '"Hedefler bölümünü SMART formatında yaz"',
                'time': '10 dk'
            },
            {
                'step': 'ADIM 5',
                'action': 'Claude Code ile Word\'e aktar',
                'prompt': '"Python-docx ile Word dökümanı oluştur"',
                'time': '3 dk'
            },
        ]

        y = Inches(2.7)
        for item in ai_steps:
            # Ana kart
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(14.4), Inches(1.1))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(25, 25, 25)
            card.fill.transparency = 0.2
            card.line.color.rgb = self.colors['success']
            card.line.width = Pt(3)

            # Step badge
            step_badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), y + Inches(0.15), Inches(1.3), Inches(0.5))
            step_badge.fill.solid()
            step_badge.fill.fore_color.rgb = self.colors['success']
            step_badge.line.fill.background()

            step_txt = s.shapes.add_textbox(Inches(1.1), y + Inches(0.15), Inches(1.3), Inches(0.5))
            step_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            step_txt.text_frame.paragraphs[0].text = item['step']
            step_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            step_txt.text_frame.paragraphs[0].font.size = Pt(20)
            step_txt.text_frame.paragraphs[0].font.bold = True
            step_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']

            # Action
            action = s.shapes.add_textbox(Inches(2.6), y + Inches(0.1), Inches(5), Inches(0.4))
            action.text_frame.paragraphs[0].text = item['action']
            action.text_frame.paragraphs[0].font.size = Pt(20)
            action.text_frame.paragraphs[0].font.bold = True
            action.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Prompt örneği
            prompt = s.shapes.add_textbox(Inches(2.6), y + Inches(0.55), Inches(10), Inches(0.4))
            prompt.text_frame.word_wrap = True
            prompt.text_frame.paragraphs[0].text = f'💬 {item["prompt"]}'
            prompt.text_frame.paragraphs[0].font.size = Pt(20)
            prompt.text_frame.paragraphs[0].font.italic = True
            prompt.text_frame.paragraphs[0].font.color.rgb = self.colors['teal']

            # Süre
            time_badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(13), y + Inches(0.3), Inches(1.5), Inches(0.5))
            time_badge.fill.solid()
            time_badge.fill.fore_color.rgb = self.colors['warning']
            time_badge.line.fill.background()

            time_txt = s.shapes.add_textbox(Inches(13), y + Inches(0.3), Inches(1.5), Inches(0.5))
            time_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            time_txt.text_frame.paragraphs[0].text = f"⏱️ {item['time']}"
            time_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            time_txt.text_frame.paragraphs[0].font.size = Pt(20)
            time_txt.text_frame.paragraphs[0].font.bold = True
            time_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']

            y += Inches(1.2)

        # 21. FAZ 2: AI ile Proposal (11:40-11:50)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'laptop', 0.3)
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(1.2), Inches(1.2))
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.colors['secondary']
        badge.line.fill.background()
        num = s.shapes.add_textbox(Inches(1), Inches(1), Inches(1.2), Inches(1.2))
        num.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        num.text_frame.paragraphs[0].text = "FAZ\n2"
        num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        num.text_frame.paragraphs[0].font.size = Pt(20)
        num.text_frame.paragraphs[0].font.bold = True
        num.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

        title = s.shapes.add_textbox(Inches(2.5), Inches(1), Inches(12), Inches(0.8))
        title.text_frame.paragraphs[0].text = "AI İLE PROPOSAL YAZIMI"
        title.text_frame.paragraphs[0].font.size = Pt(42)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

        time = s.shapes.add_textbox(Inches(2.5), Inches(1.9), Inches(8), Inches(0.5))
        time.text_frame.paragraphs[0].text = "⏰ 11:40-11:50 (10 dk)"
        time.text_frame.paragraphs[0].font.size = Pt(20)
        time.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3), Inches(13), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(25, 25, 25)
        box.fill.transparency = 0.2
        box.line.color.rgb = self.colors['secondary']
        box.line.width = Pt(3)

        steps = [
            '👉 Yukarıdaki 5 adımı takip edin',
            '👉 Claude DeepResearch ile literatür tarayın',
            '👉 TÜBİTAK 2209-A formatını kullanın',
            '👉 9 bölümü AI yardımıyla doldurun',
            '👉 Python-docx ile Word\'e aktarın (opsiyonel)',
            '👉 Takım arkadaşınızla gözden geçirin',
            '',
            '✅ ÇIKTI: 3-5 sayfa proje önerisi (Word/PDF)'
        ]

        txt = s.shapes.add_textbox(Inches(2), Inches(3.3), Inches(12), Inches(3.9))
        txt.text_frame.word_wrap = True
        for i, step in enumerate(steps):
            if i > 0:
                p = txt.text_frame.add_paragraph()
            else:
                p = txt.text_frame.paragraphs[0]
            p.text = step
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['white']
            p.space_before = Pt(20)
            if step.startswith('✅'):
                p.font.bold = True
                p.font.color.rgb = self.colors['success']

        # 22. FAZ 3: Hızlı Pitch (11:50-12:00)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'presentation', 0.3)
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(1.2), Inches(1.2))
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.colors['secondary']
        badge.line.fill.background()
        num = s.shapes.add_textbox(Inches(1), Inches(1), Inches(1.2), Inches(1.2))
        num.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        num.text_frame.paragraphs[0].text = "FAZ\n3"
        num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        num.text_frame.paragraphs[0].font.size = Pt(20)
        num.text_frame.paragraphs[0].font.bold = True
        num.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

        title = s.shapes.add_textbox(Inches(2.5), Inches(1), Inches(12), Inches(0.8))
        title.text_frame.paragraphs[0].text = "HIZLI PITCH"
        title.text_frame.paragraphs[0].font.size = Pt(42)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

        time = s.shapes.add_textbox(Inches(2.5), Inches(1.9), Inches(8), Inches(0.5))
        time.text_frame.paragraphs[0].text = "⏰ 11:50-12:00 (10 dk)"
        time.text_frame.paragraphs[0].font.size = Pt(20)
        time.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']

        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3), Inches(13), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(25, 25, 25)
        box.fill.transparency = 0.2
        box.line.color.rgb = self.colors['secondary']
        box.line.width = Pt(3)

        steps = [
            '👉 1 dakikalık pitch metni yazın',
            '👉 AI ile düzeltin ve geliştirin',
            '👉 Takım içinde prova yapın (30 sn)',
            '👉 Gönüllü 3-5 takım sunar',
            '',
            '🎤 PITCH FORMATI:',
            '   • Problem (15 sn): "Fabrikalar enerji israf ediyor"',
            '   • Çözüm (25 sn): "IoT + AI ile %30 tasarruf"',
            '   • Etki (20 sn): "Yılda 2M TL tasarruf, 100 ton CO2"',
            '',
            '🏆 EN İYİ 3 PROJE ÖDÜL KAZANIR!'
        ]

        txt = s.shapes.add_textbox(Inches(2), Inches(3.3), Inches(12), Inches(3.9))
        txt.text_frame.word_wrap = True
        for i, step in enumerate(steps):
            if i > 0:
                p = txt.text_frame.add_paragraph()
            else:
                p = txt.text_frame.paragraphs[0]
            p.text = step
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['white']
            p.space_before = Pt(6)
            if step.startswith('🏆') or step.startswith('🎤'):
                p.font.bold = True
                p.font.color.rgb = self.colors['warning'] if step.startswith('🎤') else self.colors['success']

        # ===== KAPANIŞ =====
        print("\n📍 KAPANIŞ...")

        # 23. Başarılar
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'success', 0.3)
        self.add_title(s, "🎯 BUGÜNÜN BAŞARILARI", 2, 54)

        achievements = [
            {'n': '25+', 'l': 'Proje\nFikri'},
            {'n': '50', 'l': 'Genç\nGirişimci'},
            {'n': '8', 'l': 'Araç\nÖğrenildi'},
            {'n': '100%', 'l': 'Motivasyon'},
        ]

        x = Inches(1.5)
        for ach in achievements:
            c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4), Inches(3), Inches(3))
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(30, 30, 30)
            c.fill.transparency = 0.2
            c.line.color.rgb = self.colors['success']
            c.line.width = Pt(3)

            nb = s.shapes.add_textbox(x, Inches(4.5), Inches(3), Inches(1.2))
            nb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            nb.text_frame.paragraphs[0].text = ach['n']
            nb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            nb.text_frame.paragraphs[0].font.size = Pt(52)
            nb.text_frame.paragraphs[0].font.bold = True
            nb.text_frame.paragraphs[0].font.color.rgb = self.colors['success']

            lb = s.shapes.add_textbox(x, Inches(5.8), Inches(3), Inches(1))
            lb.text_frame.word_wrap = True
            lb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            lb.text_frame.paragraphs[0].text = ach['l']
            lb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            lb.text_frame.paragraphs[0].font.size = Pt(20)
            lb.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            x += Inches(3.3)

        # 24. Sonraki Adımlar (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'next_steps', 0.35)
        self.add_title(s, "🚀 SONRAKİ ADIMLAR", 1, 44)

        next_steps = [
            {
                'timeline': 'Bu Hafta',
                'tasks': [
                    {'task': 'Proje önerinizi tamamlayın', 'detail': 'AI ile son düzeltmeleri yapın'},
                    {'task': 'Takım arkadaşlarınızla gözden geçirin', 'detail': '3-2-1 feedback metoduyla'},
                    {'task': 'Danışman öğretim üyesi bulun', 'detail': 'TÜBİTAK için zorunlu'},
                ],
                'color': self.colors['danger']
            },
            {
                'timeline': 'Bu Ay (Kasım)',
                'tasks': [
                    {'task': 'Şirket kuruluşu (KOSGEB için)', 'detail': 'Noter + Ticaret Sicili'},
                    {'task': 'KOSGEB Girişimcilik eğitimi', 'detail': '40 saat online eğitim'},
                    {'task': 'Prototip/MVP geliştirmeye başlayın', 'detail': 'Claude Code ile hızlı ilerleme'},
                ],
                'color': self.colors['warning']
            },
            {
                'timeline': 'Aralık-Ocak',
                'tasks': [
                    {'task': 'TÜBİTAK 2209-A başvurusu', 'detail': 'Şubat deadline (kesin tarihi web\'ten kontrol)'},
                    {'task': 'Patent/Faydalı model araştırması', 'detail': 'Türk Patent Enstitüsü'},
                    {'task': 'İlk pilot testler', 'detail': '5-10 kullanıcı ile beta test'},
                ],
                'color': self.colors['primary']
            },
            {
                'timeline': 'Şubat-Mart',
                'tasks': [
                    {'task': 'Teknofest başvuruları', 'detail': '40+ kategori, ödüller 100K TL\'ye kadar'},
                    {'task': 'KOSGEB 1. dönem başvurusu', 'detail': 'İş planı + finansal tablolar hazır olsun'},
                    {'task': 'Pitch deck hazırlayın', 'detail': 'Yatırımcı sunumları için'},
                ],
                'color': self.colors['success']
            },
        ]

        y = Inches(2.5)  # CHANGED: 3.0"->2.5" daha yukarıdan başla
        for step in next_steps:
            # Timeline header
            timeline_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), y, Inches(3.2), Inches(0.55))  # CHANGED: daha geniş header
            timeline_box.fill.solid()
            timeline_box.fill.fore_color.rgb = step['color']
            timeline_box.line.fill.background()

            timeline_txt = s.shapes.add_textbox(Inches(1.2), y, Inches(3.2), Inches(0.55))
            timeline_txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            timeline_txt.text_frame.paragraphs[0].text = step['timeline']
            timeline_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            timeline_txt.text_frame.paragraphs[0].font.size = Pt(22)  # CHANGED: 20->22 daha belirgin
            timeline_txt.text_frame.paragraphs[0].font.bold = True
            timeline_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

            # Tasks card - COMPLETELY REORGANIZED
            task_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), y, Inches(10.1), Inches(1.05))  # CHANGED: daha geniş ve yüksek kart
            task_card.fill.solid()
            task_card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            task_card.fill.transparency = 0.25
            task_card.line.color.rgb = step['color']
            task_card.line.width = Pt(2)

            # Tasks list - REORGANIZED with better spacing
            task_y = y + Inches(0.12)  # CHANGED: daha fazla üst padding
            for i, task in enumerate(step['tasks']):
                task_txt = s.shapes.add_textbox(Inches(5.0), task_y, Inches(9.5), Inches(0.28))  # CHANGED: daha geniş ve yüksek textbox
                task_txt.text_frame.word_wrap = True
                task_txt.text_frame.paragraphs[0].text = f"• {task['task']}: {task['detail']}"
                task_txt.text_frame.paragraphs[0].font.size = Pt(18)  # CHANGED: 20->18 daha kompakt
                task_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
                task_y += Inches(0.30)  # CHANGED: 0.27"->0.30" daha fazla spacing

            y += Inches(1.40)  # CHANGED: 1.05"->1.40" daha fazla kart arası boşluk, total: 2.5+(4×1.40)=8.1" ✓

        # 25. Kaynaklar ve Linkler (YENİ)
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'resources', 0.35)
        self.add_title(s, "📚 KAYNAKLAR VE LİNKLER", 1, 44)

        resources = [
            {
                'category': '🤖 AI Araçları',
                'items': [
                    {'name': 'Claude.ai', 'link': 'claude.ai', 'desc': 'En iyi araştırma & yazma asistanı'},
                    {'name': 'Claude Code', 'link': 'claude.com/claude-code', 'desc': 'Terminal AI yardımcısı'},
                    {'name': 'ChatGPT', 'link': 'chat.openai.com', 'desc': 'Genel amaçlı AI'},
                    {'name': 'Google Gemini', 'link': 'gemini.google.com', 'desc': 'Google\'ın AI modeli'},
                ]
            },
            {
                'category': '💰 Fonlama',
                'items': [
                    {'name': 'TÜBİTAK 2209-A', 'link': 'tubitak.gov.tr/tr/burslar/lisans-onlisans/destek-programlari/2209-universite-ogrencileri-arastirma-projeleri-destekleme-programi', 'desc': 'Öğrenci projeleri'},
                    {'name': 'KOSGEB', 'link': 'kosgeb.gov.tr', 'desc': 'Girişim destekleri'},
                    {'name': 'Teknofest', 'link': 'teknofest.org/tr', 'desc': 'Teknoloji yarışmaları'},
                    {'name': 'Horizon Europe', 'link': 'ec.europa.eu/horizon-europe', 'desc': 'AB fonları'},
                ]
            },
            {
                'category': '📖 Öğrenme',
                'items': [
                    {'name': 'Google Scholar', 'link': 'scholar.google.com', 'desc': 'Akademik makaleler'},
                    {'name': 'Coursera', 'link': 'coursera.org', 'desc': 'Ücretsiz online kurslar'},
                    {'name': 'MIT OpenCourseWare', 'link': 'ocw.mit.edu', 'desc': 'MIT dersleri'},
                    {'name': 'Kaggle', 'link': 'kaggle.com', 'desc': 'Veri bilimi & ML'},
                ]
            },
            {
                'category': '🛠️ Araçlar',
                'items': [
                    {'name': 'GitHub', 'link': 'github.com', 'desc': 'Kod depolama & işbirliği'},
                    {'name': 'Notion', 'link': 'notion.so', 'desc': 'Proje yönetimi'},
                    {'name': 'Figma', 'link': 'figma.com', 'desc': 'UI/UX tasarım'},
                    {'name': 'Loom', 'link': 'loom.com', 'desc': 'Ekran kaydı & demo'},
                ]
            },
        ]

        cols = 2
        rows = 2
        x_start = 1.2
        y_start = 2.2  # CHANGED: daha yukarı
        w = 6.8  # CHANGED: daha geniş kartlar
        h = 3.0  # CHANGED: daha kompakt kartlar
        x_gap = 0.4
        y_gap = 0.3

        for i, resource in enumerate(resources):
            row = i // cols
            col = i % cols
            x = x_start + col * (w + x_gap)
            y = y_start + row * (h + y_gap)

            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(20, 20, 20)
            card.fill.transparency = 0.25
            card.line.color.rgb = self.colors['teal']
            card.line.width = Pt(3)

            # Category header
            cat_txt = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(w - 0.4), Inches(0.35))
            cat_txt.text_frame.paragraphs[0].text = resource['category']
            cat_txt.text_frame.paragraphs[0].font.size = Pt(22)  # CHANGED: 20->22 daha belirgin
            cat_txt.text_frame.paragraphs[0].font.bold = True
            cat_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['warning']
            cat_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Items - COMPLETELY REORGANIZED
            item_y = y + 0.58
            for item in resource['items']:
                # Name (bold)
                name_txt = s.shapes.add_textbox(Inches(x + 0.25), Inches(item_y), Inches(w - 0.5), Inches(0.20))
                name_txt.text_frame.paragraphs[0].text = f"• {item['name']}"
                name_txt.text_frame.paragraphs[0].font.size = Pt(18)  # CHANGED: 20->18 daha kompakt
                name_txt.text_frame.paragraphs[0].font.bold = True
                name_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

                # Link (readable size!)
                link_txt = s.shapes.add_textbox(Inches(x + 0.4), Inches(item_y + 0.20), Inches(w - 0.65), Inches(0.18))
                link_txt.text_frame.word_wrap = True
                link_txt.text_frame.paragraphs[0].text = f"🔗 {item['link']}"
                link_txt.text_frame.paragraphs[0].font.size = Pt(14)  # CHANGED: 6->14 READABLE!
                link_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

                # Description
                desc_txt = s.shapes.add_textbox(Inches(x + 0.4), Inches(item_y + 0.38), Inches(w - 0.65), Inches(0.16))
                desc_txt.text_frame.word_wrap = True
                desc_txt.text_frame.paragraphs[0].text = item['desc']
                desc_txt.text_frame.paragraphs[0].font.size = Pt(16)  # CHANGED: 20->16 daha kompakt
                desc_txt.text_frame.paragraphs[0].font.color.rgb = self.colors['light']

                item_y += 0.60  # CHANGED: 0.50->0.60 daha fazla boşluk

        # 26. İletişim
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_counter += 1
        self.add_bg_image(s, 'rocket', 0.5)
        self.add_title(s, "GELECEĞİ BİRLİKTE İNŞA EDELİM!", 2.8, 48)

        contact = s.shapes.add_textbox(Inches(4), Inches(4.5), Inches(8), Inches(3))
        lines = [
            "Dr. Öğr. Üyesi Uğur CORUH",
            "",
            "📧 ugur.coruh@erdogan.edu.tr",
            "💼 LinkedIn: linkedin.com/in/ugurcoruh",
            "💻 GitHub: github.com/ucoruh",
            "",
            "RTEÜ Erasmus+"
        ]
        for i, line in enumerate(lines):
            if i > 0:
                p = contact.text_frame.add_paragraph()
            else:
                p = contact.text_frame.paragraphs[0]
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['white']
            if i == 0:
                p.font.size = Pt(26)
                p.font.bold = True

        # Kaydet
        import datetime
        timestamp = datetime.datetime.now().strftime("%H%M")
        fn = f'RTEU_Workshop_{timestamp}.pptx'
        self.prs.save(fn)

        print("\n" + "="*70)
        print(f"✅ KAPSAMLI SUNUM OLUŞTURULDU: {fn}")
        print("="*70)
        print(f"📊 Toplam slayt: {len(self.prs.slides)}")
        print(f"\n📸 Görsel istatistikleri:")
        print(f"   • Toplam: {self.stats['total']}")
        print(f"   • Başarılı: {self.stats['success']} ✅")
        print(f"   • Başarısız: {self.stats['failed']} ❌")
        if self.stats['total'] > 0:
            print(f"   • Başarı: %{(self.stats['success']/self.stats['total'])*100:.1f}")

        print(f"\n🎯 SUNUM YAPISI:")
        print(f"   📍 Açılış: 6 slayt (GENİŞLETİLDİ!)")
        print(f"   📍 Sanayi 4.0: 10 slayt (GENİŞLETİLDİ!)")
        print(f"   📍 AI Araçları: 12 slayt (GENİŞLETİLDİ!)")
        print(f"   📍 Fonlama: 8 slayt (GENİŞLETİLDİ!)")
        print(f"   📍 Mola & Hazırlık: 2 slayt (GENİŞLETİLDİ!)")
        print(f"   📍 Atölye: 18 slayt (GENİŞLETİLDİ!)")
        print(f"   📍 Kapanış: 4 slayt (GENİŞLETİLDİ!)")
        print("="*70 + "\n")

if __name__ == "__main__":
    print("="*70)
    print("🎯 ULTİMATE 2 SAATLİK KAPSAMLI WORKSHOP SUNUMU")
    print("="*70)

    try:
        p = UltimatePresentation()
        p.generate()
        print("✅ Detaylı sunum hazır! Proje yazım formatları ve")
        print("   adım adım yönergelerle öğrenciler ne yapacağını biliyor!")
    except Exception as e:
        print(f"❌ Hata: {e}")
        import traceback
        traceback.print_exc()
